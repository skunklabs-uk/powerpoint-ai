#!/usr/bin/env python3
"""Apply a single coverage axis to the user-refined comparison slide."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Data pipeline blueprint v3.pptx"
OUTPUT = ROOT / "attempts" / "Data pipeline blueprint v5 - coverage definitions.pptx"

FONT = "Poppins"
MUTED = RGBColor(100, 113, 122)

STYLES = {
    "direct": {
        "label": "DIRETTA",
        "fill": RGBColor(234, 248, 240),
        "accent": RGBColor(51, 175, 116),
    },
    "config": {
        "label": "CONFIGURAZIONE",
        "fill": RGBColor(239, 248, 252),
        "accent": RGBColor(34, 153, 211),
    },
    "extension": {
        "label": "ESTENSIONE",
        "fill": RGBColor(253, 247, 233),
        "accent": RGBColor(220, 160, 55),
    },
    "uncovered": {
        "label": "NON COPERTO",
        "fill": RGBColor(242, 244, 245),
        "accent": RGBColor(128, 140, 148),
    },
}

# Top-level group indexes on slide 11, preserved from the user-refined layout.
AWS_GROUPS = [8, 11, 14, 17, 20, 23, 26, 29]
QLIK_GROUPS = [9, 12, 15, 18, 21, 24, 27, 30]

AWS_CLASSIFICATION = [
    ("extension", False),
    ("direct", False),
    ("direct", False),
    ("extension", False),
    ("config", False),
    ("extension", False),
    ("direct", False),
    ("config", False),
]

QLIK_CLASSIFICATION = [
    ("direct", False),
    ("extension", True),
    ("direct", False),
    ("extension", False),
    ("config", False),
    ("extension", False),
    ("direct", False),
    ("config", True),
]

STATUS_LABELS = {
    "COPERTO",
    "DA CONFIGURARE",
    "DA COMPLETARE",
    "DA VERIFICARE",
}


def set_single_run_text(text_box, text, color=None):
    paragraph = text_box.text_frame.paragraphs[0]
    if not paragraph.runs:
        run = paragraph.add_run()
    else:
        run = paragraph.runs[0]
        for extra in paragraph.runs[1:]:
            extra._r.getparent().remove(extra._r)
    run.text = text
    if color is not None:
        run.font.color.rgb = color


def find_status_parts(group_shape):
    shapes = group_shape.shapes
    for index, shape in enumerate(shapes):
        text = " ".join(shape.text.split()) if hasattr(shape, "text") else ""
        if text in STATUS_LABELS and index > 0:
            return shapes[index - 1], shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                return find_status_parts(shape)
            except LookupError:
                pass
    raise LookupError(f"Status label not found in group: {group_shape.name}")


def set_status(slide, top_group_index, kind, uncertain):
    style = STYLES[kind]
    pill, label = find_status_parts(slide.shapes[top_group_index])

    pill.fill.solid()
    pill.fill.fore_color.rgb = style["fill"]
    pill.line.color.rgb = style["accent"]
    pill.line.width = Pt(0.8)

    text = style["label"] + ("*" if uncertain else "")
    set_single_run_text(label, text, style["accent"])


def update_title(slide):
    title = slide.shapes[0]
    paragraph = title.text_frame.paragraphs[0]
    if len(paragraph.runs) >= 2:
        paragraph.runs[1].text = "Modalità di copertura di ciascuna soluzione"
    else:
        title.text = "Data pipeline | Modalità di copertura di ciascuna soluzione"


def update_legend(slide):
    old_legend = slide.shapes[31]
    old_legend.element.getparent().remove(old_legend.element)

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.76), Inches(6.39), Inches(11.48), Inches(0.008)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(218, 228, 233)
    divider.line.fill.background()

    add_text(
        slide,
        0.76,
        6.44,
        0.78,
        0.18,
        "LEGENDA",
        7.6,
        MUTED,
        True,
    )
    add_text(
        slide,
        8.52,
        6.24,
        3.72,
        0.16,
        "* capacità da verificare: edizione, licenza o componente da confermare",
        6.3,
        STYLES["extension"]["accent"],
        True,
        PP_ALIGN.RIGHT,
    )

    definitions = [
        (
            "DIRETTA",
            "Funzione disponibile con i componenti previsti",
            STYLES["direct"]["accent"],
        ),
        (
            "CONFIGURAZIONE",
            "Richiede parametri, regole o impostazioni",
            STYLES["config"]["accent"],
        ),
        (
            "ESTENSIONE",
            "Richiede sviluppo, connettori o integrazioni aggiuntive",
            STYLES["extension"]["accent"],
        ),
        (
            "NON COPERTO",
            "Funzione esterna al perimetro della soluzione",
            STYLES["uncovered"]["accent"],
        ),
    ]

    start_x = 1.62
    item_width = 2.64
    for index, (label, description, color) in enumerate(definitions):
        x = start_x + index * item_width
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(6.49), Inches(0.11), Inches(0.11)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text(slide, x + 0.17, 6.43, 2.36, 0.18, label, 7.3, color, True)
        add_text(slide, x + 0.17, 6.61, 2.33, 0.34, description, 6.7, MUTED)

        if index < len(definitions) - 1:
            separator = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + 2.54),
                Inches(6.46),
                Inches(0.006),
                Inches(0.49),
            )
            separator.fill.solid()
            separator.fill.fore_color.rgb = RGBColor(226, 233, 237)
            separator.line.fill.background()


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def build():
    if OUTPUT.exists():
        raise FileExistsError(f"Refusing to overwrite existing output: {OUTPUT}")

    prs = Presentation(SOURCE)
    slide = prs.slides[10]
    update_title(slide)

    for group_index, (kind, uncertain) in zip(AWS_GROUPS, AWS_CLASSIFICATION):
        set_status(slide, group_index, kind, uncertain)
    for group_index, (kind, uncertain) in zip(QLIK_GROUPS, QLIK_CLASSIFICATION):
        set_status(slide, group_index, kind, uncertain)

    update_legend(slide)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
