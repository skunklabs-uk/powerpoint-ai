#!/usr/bin/env python3
"""Add the executive process-coverage comparison slide to the v2 deck."""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Data pipeline blueprint v2.pptx"
OUTPUT = ROOT / "Data pipeline blueprint v3.pptx"

FONT = "Poppins"
INK = RGBColor(44, 55, 64)
MUTED = RGBColor(100, 113, 122)
LINE = RGBColor(218, 228, 233)
LIGHT_ROW = RGBColor(248, 250, 251)
BLUE = RGBColor(34, 153, 211)
BLUE_LIGHT = RGBColor(239, 248, 252)
TEAL = RGBColor(38, 177, 172)
TEAL_LIGHT = RGBColor(239, 250, 248)
GREEN = RGBColor(51, 175, 116)
GREEN_LIGHT = RGBColor(234, 248, 240)
AMBER = RGBColor(220, 160, 55)
AMBER_LIGHT = RGBColor(253, 247, 233)


ROWS = [
    (
        "Fonti e acquisizione",
        ("DA CONFIGURARE", "Connettori dedicati, coordinati da Dagster", "config"),
        ("COPERTO", "Acquisizione con Qlik Talend Cloud", "covered"),
    ),
    (
        "Conservazione degli input",
        ("COPERTO", "Database e archivio file su AWS", "covered"),
        ("DA COMPLETARE", "Archivio storico ancora da scegliere", "attention"),
    ),
    (
        "Preparazione e mapping",
        ("COPERTO", "dbt e tabelle di corrispondenza", "covered"),
        ("COPERTO", "Preparazione e mapping con Talend", "covered"),
    ),
    (
        "Regole, allocazioni e calcoli",
        ("DA CONFIGURARE", "Modelli dbt e regole CDG da implementare", "config"),
        ("DA CONFIGURARE", "Procedure Talend, SQL o moduli dedicati", "config"),
    ),
    (
        "Qualità e riconciliazione",
        ("DA CONFIGURARE", "Test, controlli e soglie da impostare", "config"),
        ("DA CONFIGURARE", "Controlli Talend e soglie da impostare", "config"),
    ),
    (
        "Actual, Forecast e versioni",
        ("DA CONFIGURARE", "Modelli separati per ciascun progetto", "config"),
        ("DA CONFIGURARE", "Procedure separate per ciascun progetto", "config"),
    ),
    (
        "Reporting e distribuzione",
        ("COPERTO", "Metabase sui dati già controllati", "covered"),
        ("COPERTO", "Qlik Cloud Analytics", "covered"),
    ),
    (
        "Monitoraggio, audit e lineage",
        ("DA CONFIGURARE", "Log, catalogo e prove delle esecuzioni", "config"),
        ("DA VERIFICARE", "Dipende dall'edizione e dalle parti personalizzate", "attention"),
    ),
]


def style_text_box(shape, *, margin=0.05):
    shape.text_frame.margin_left = Inches(margin)
    shape.text_frame.margin_right = Inches(margin)
    shape.text_frame.margin_top = Inches(0.01)
    shape.text_frame.margin_bottom = Inches(0.01)
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text(slide, x, y, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    style_text_box(box)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_round_rect(slide, x, y, w, h, fill, line=LINE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_header_card(slide, x, width, accent, fill, label, subtitle):
    add_round_rect(slide, x, 1.43, width, 0.68, fill, accent)
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.43), Inches(0.08), Inches(0.68)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    add_text(slide, x + 0.22, 1.50, width - 0.40, 0.22, label, 13.5, INK, True)
    add_text(slide, x + 0.22, 1.75, width - 0.40, 0.22, subtitle, 8.3, MUTED)


def status_colors(kind):
    if kind == "covered":
        return GREEN_LIGHT, GREEN
    if kind == "config":
        return BLUE_LIGHT, BLUE
    return AMBER_LIGHT, AMBER


def add_status_cell(slide, x, y, width, status, description, kind, alternate):
    cell_fill = LIGHT_ROW if alternate else RGBColor(255, 255, 255)
    add_round_rect(slide, x, y, width, 0.41, cell_fill, LINE)

    pill_fill, accent = status_colors(kind)
    pill = add_round_rect(slide, x + 0.10, y + 0.075, 1.20, 0.26, pill_fill, accent)
    pill.line.width = Pt(0.7)
    add_text(slide, x + 0.13, y + 0.086, 1.14, 0.22, status, 6.8, accent, True, PP_ALIGN.CENTER)
    add_text(slide, x + 1.42, y + 0.055, width - 1.53, 0.30, description, 8.15, INK)


def add_process_cell(slide, x, y, width, index, label, alternate):
    cell_fill = LIGHT_ROW if alternate else RGBColor(255, 255, 255)
    add_round_rect(slide, x, y, width, 0.41, cell_fill, LINE)

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.10), Inches(y + 0.075), Inches(0.26), Inches(0.26)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    circle.line.color.rgb = TEAL
    circle.line.width = Pt(1.0)
    add_text(slide, x + 0.10, y + 0.082, 0.26, 0.22, str(index), 7.0, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.46, y + 0.055, width - 0.56, 0.30, label, 8.7, INK, True)


def replace_title(slide, source_title):
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape.element.getparent().remove(shape.element)

    title_xml = deepcopy(source_title.element)
    text_nodes = title_xml.xpath(".//a:t")
    text_nodes[0].text = "Data pipeline | "
    text_nodes[1].text = "Il processo è coperto da entrambe; cambia il livello di integrazione"
    slide.shapes._spTree.insert_element_before(title_xml, "p:extLst")


def build():
    if OUTPUT.exists():
        raise FileExistsError(f"Refusing to overwrite existing output: {OUTPUT}")

    prs = Presentation(SOURCE)
    template_slide = prs.slides[12]
    slide = prs.slides.add_slide(template_slide.slide_layout)
    replace_title(slide, template_slide.shapes[0])

    # Move the new slide after "Un modello, due modi per realizzarla".
    slide_id = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(10, slide_id)

    band = add_round_rect(slide, 0.78, 0.79, 11.57, 0.48, TEAL_LIGHT, TEAL)
    band.line.width = Pt(0.8)
    add_text(
        slide,
        0.98,
        0.865,
        11.17,
        0.30,
        "Copertura tecnologica non significa progetto pronto: fonti, regole, soglie e responsabilità restano da definire in entrambi gli scenari.",
        9.4,
        INK,
        False,
        PP_ALIGN.CENTER,
    )

    process_x, process_w = 0.78, 2.65
    aws_x, solution_w = 3.61, 4.28
    qlik_x = 8.07

    add_round_rect(slide, process_x, 1.43, process_w, 0.68, RGBColor(246, 248, 249), LINE)
    add_text(slide, process_x + 0.20, 1.53, process_w - 0.40, 0.19, "PROCESSO END-TO-END", 9.2, MUTED, True)
    add_text(slide, process_x + 0.20, 1.77, process_w - 0.40, 0.19, "8 passaggi della blueprint", 8.1, MUTED)

    add_header_card(
        slide,
        aws_x,
        solution_w,
        BLUE,
        BLUE_LIGHT,
        "AWS | Stack componibile",
        "8/8 passaggi copribili con componenti coordinati",
    )
    add_header_card(
        slide,
        qlik_x,
        solution_w,
        TEAL,
        TEAL_LIGHT,
        "Qlik | Piattaforma Qlik + Talend",
        "8/8 passaggi copribili in un ambiente più integrato",
    )

    start_y = 2.22
    row_step = 0.47
    for index, (label, aws, qlik) in enumerate(ROWS, start=1):
        y = start_y + (index - 1) * row_step
        alternate = index % 2 == 0
        add_process_cell(slide, process_x, y, process_w, index, label, alternate)
        add_status_cell(slide, aws_x, y, solution_w, *aws, alternate)
        add_status_cell(slide, qlik_x, y, solution_w, *qlik, alternate)

    legend_y = 6.08
    add_text(slide, 0.82, legend_y + 0.02, 1.12, 0.20, "LEGENDA", 8.0, MUTED, True)
    legend = [
        (GREEN, "Coperto dalla soluzione"),
        (BLUE, "Configurazione o sviluppo mirato"),
        (AMBER, "Da completare o verificare"),
    ]
    lx = 1.70
    for color, label in legend:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx), Inches(legend_y + 0.04), Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text(slide, lx + 0.20, legend_y, 2.48, 0.22, label, 7.7, MUTED)
        lx += 3.02

    add_round_rect(slide, 0.78, 6.36, 5.65, 0.46, BLUE_LIGHT, BLUE)
    add_text(
        slide,
        0.99,
        6.40,
        5.24,
        0.31,
        "AWS: maggiore modularità e sostituibilità; maggiore onere di integrazione e gestione.",
        8.4,
        INK,
        True,
        PP_ALIGN.CENTER,
    )
    add_round_rect(slide, 6.70, 6.36, 5.65, 0.46, TEAL_LIGHT, TEAL)
    add_text(
        slide,
        6.91,
        6.40,
        5.24,
        0.31,
        "Qlik: maggiore integrazione e rapidità potenziale; maggiore dipendenza da edizione, licenze e piattaforma.",
        8.4,
        INK,
        True,
        PP_ALIGN.CENTER,
    )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
