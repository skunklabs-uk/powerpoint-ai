#!/usr/bin/env python3
"""Add project-presence tags to the process coverage slide."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "attempts" / "Data pipeline blueprint v5 - coverage definitions.pptx"
OUTPUT = ROOT / "attempts" / "Data pipeline blueprint v9 - project process tags.pptx"

FONT = "Poppins"
INK = RGBColor(44, 55, 64)
MUTED = RGBColor(100, 113, 122)

PROCESS_GROUPS = [7, 10, 13, 16, 19, 22, 25, 28]

PROJECTS = {
    "ProSIGNAL": {
        "label": "ProSIGNAL",
        "width": 0.59,
        "fill": RGBColor(239, 248, 252),
        "accent": RGBColor(34, 153, 211),
    },
    "Kiron CDG": {
        "label": "Kiron CDG",
        "width": 0.61,
        "fill": RGBColor(239, 250, 248),
        "accent": RGBColor(38, 177, 172),
    },
    "CDG interno": {
        "label": "CDG interno",
        "width": 0.68,
        "fill": RGBColor(234, 248, 240),
        "accent": RGBColor(51, 175, 116),
    },
}

# Presence means that the process is explicitly documented in the available
# project material. It does not imply implementation or production readiness.
PROCESS_PROJECTS = [
    ("ProSIGNAL", "Kiron CDG", "CDG interno"),
    ("ProSIGNAL", "Kiron CDG", "CDG interno"),
    ("ProSIGNAL", "Kiron CDG", "CDG interno"),
    ("ProSIGNAL", "Kiron CDG", "CDG interno"),
    ("ProSIGNAL", "Kiron CDG", "CDG interno"),
    ("Kiron CDG", "CDG interno"),
    ("ProSIGNAL", "Kiron CDG", "CDG interno"),
    ("Kiron CDG", "CDG interno"),
]


def set_text(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, color, bold, align)
    return shape


def update_title(slide):
    title = slide.shapes[0]
    paragraph = title.text_frame.paragraphs[0]
    if len(paragraph.runs) >= 2:
        paragraph.runs[1].text = "Processi condivisi dai progetti e modalità di copertura"
    else:
        set_text(
            title,
            "Data pipeline | Processi condivisi dai progetti e modalità di copertura",
            17.5,
            INK,
            True,
        )


def update_process_header(slide):
    header = slide.shapes[2]
    header.top = Inches(1.49)
    header.height = Inches(0.19)
    set_text(header, "PROCESSO / PROGETTI", 8.7, MUTED, True)
    add_text(
        slide,
        0.98,
        1.75,
        2.16,
        0.16,
        "Presenza nei casi analizzati",
        7.0,
        MUTED,
    )


def add_project_tag(slide, x, y, project):
    style = PROJECTS[project]
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(style["width"]),
        Inches(0.135),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = style["fill"]
    tag.line.color.rgb = style["accent"]
    tag.line.width = Pt(0.55)
    set_text(tag, style["label"], 5.1, style["accent"], True, PP_ALIGN.CENTER)
    return style["width"]


def update_process_cells(slide):
    for group_index, projects in zip(PROCESS_GROUPS, PROCESS_PROJECTS):
        group = slide.shapes[group_index]
        label = group.shapes[3]
        label.top = group.top + Inches(0.025)
        label.height = Inches(0.18)
        set_text(label, label.text, 7.65, INK, True)

        x = 1.24
        y = group.top / 914400 + 0.235
        for project in projects:
            x += add_project_tag(slide, x, y, project) + 0.045


def add_presence_note(slide):
    add_text(
        slide,
        0.76,
        6.24,
        4.72,
        0.15,
        "Tag progetto = processo documentato; non indica che sia già implementato",
        6.3,
        MUTED,
        True,
    )


def clear_zero_size_master_labels(slide):
    master = slide.slide_layout.slide_master
    for shape in master.shapes:
        if (shape.width == 0 or shape.height == 0) and shape.has_text_frame:
            shape.text = ""


def build():
    if OUTPUT.exists():
        raise FileExistsError(f"Refusing to overwrite existing output: {OUTPUT}")

    prs = Presentation(SOURCE)
    slide = prs.slides[10]
    clear_zero_size_master_labels(slide)
    update_title(slide)
    update_process_header(slide)
    update_process_cells(slide)
    add_presence_note(slide)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
