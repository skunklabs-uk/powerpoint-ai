#!/usr/bin/env python3
"""Apply a single coverage axis to the user-refined comparison slide."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Data pipeline blueprint v3.pptx"
OUTPUT = ROOT / "attempts" / "Data pipeline blueprint v4 - coverage axis.pptx"

FONT = "Poppins"
MUTED = RGBColor(100, 113, 122)

STYLES = {
    "direct": {
        "label": "DIRETTA",
        "fill": RGBColor(234, 248, 240),
        "accent": RGBColor(51, 175, 116),
    },
    "config": {
        "label": "CONFIGURAZIONE",
        "fill": RGBColor(239, 248, 252),
        "accent": RGBColor(34, 153, 211),
    },
    "extension": {
        "label": "ESTENSIONE",
        "fill": RGBColor(253, 247, 233),
        "accent": RGBColor(220, 160, 55),
    },
}

# Top-level group indexes on slide 11, preserved from the user-refined layout.
AWS_GROUPS = [8, 11, 14, 17, 20, 23, 26, 29]
QLIK_GROUPS = [9, 12, 15, 18, 21, 24, 27, 30]

AWS_CLASSIFICATION = [
    ("extension", False),
    ("direct", False),
    ("direct", False),
    ("extension", False),
    ("config", False),
    ("extension", False),
    ("direct", False),
    ("config", False),
]

QLIK_CLASSIFICATION = [
    ("direct", False),
    ("extension", True),
    ("direct", False),
    ("extension", False),
    ("config", False),
    ("extension", False),
    ("direct", False),
    ("config", True),
]

STATUS_LABELS = {
    "COPERTO",
    "DA CONFIGURARE",
    "DA COMPLETARE",
    "DA VERIFICARE",
}


def set_single_run_text(text_box, text, color=None):
    paragraph = text_box.text_frame.paragraphs[0]
    if not paragraph.runs:
        run = paragraph.add_run()
    else:
        run = paragraph.runs[0]
        for extra in paragraph.runs[1:]:
            extra._r.getparent().remove(extra._r)
    run.text = text
    if color is not None:
        run.font.color.rgb = color


def find_status_parts(group_shape):
    shapes = group_shape.shapes
    for index, shape in enumerate(shapes):
        text = " ".join(shape.text.split()) if hasattr(shape, "text") else ""
        if text in STATUS_LABELS and index > 0:
            return shapes[index - 1], shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                return find_status_parts(shape)
            except LookupError:
                pass
    raise LookupError(f"Status label not found in group: {group_shape.name}")


def set_status(slide, top_group_index, kind, uncertain):
    style = STYLES[kind]
    pill, label = find_status_parts(slide.shapes[top_group_index])

    pill.fill.solid()
    pill.fill.fore_color.rgb = style["fill"]
    pill.line.color.rgb = style["accent"]
    pill.line.width = Pt(0.8)

    text = style["label"] + ("*" if uncertain else "")
    set_single_run_text(label, text, style["accent"])


def update_title(slide):
    title = slide.shapes[0]
    paragraph = title.text_frame.paragraphs[0]
    if len(paragraph.runs) >= 2:
        paragraph.runs[1].text = "Modalità di copertura di ciascuna soluzione"
    else:
        title.text = "Data pipeline | Modalità di copertura di ciascuna soluzione"


def update_legend(slide):
    legend = slide.shapes[31].shapes[1]
    labels = [
        ("Copertura diretta", STYLES["direct"]["accent"]),
        ("Configurazione", STYLES["config"]["accent"]),
        ("Estensione / integrazione", STYLES["extension"]["accent"]),
    ]

    for group, (text, color) in zip(legend.shapes, labels):
        dot = group.shapes[0]
        text_box = group.shapes[1]
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        set_single_run_text(text_box, text, MUTED)

    note = slide.shapes.add_textbox(
        Inches(10.42), Inches(6.735), Inches(1.95), Inches(0.22)
    )
    note.text_frame.clear()
    note.text_frame.margin_left = 0
    note.text_frame.margin_right = 0
    note.text_frame.margin_top = 0
    note.text_frame.margin_bottom = 0
    note.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = note.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = "* capacità da verificare"
    run.font.name = FONT
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.color.rgb = STYLES["extension"]["accent"]


def build():
    if OUTPUT.exists():
        raise FileExistsError(f"Refusing to overwrite existing output: {OUTPUT}")

    prs = Presentation(SOURCE)
    slide = prs.slides[10]
    update_title(slide)

    for group_index, (kind, uncertain) in zip(AWS_GROUPS, AWS_CLASSIFICATION):
        set_status(slide, group_index, kind, uncertain)
    for group_index, (kind, uncertain) in zip(QLIK_GROUPS, QLIK_CLASSIFICATION):
        set_status(slide, group_index, kind, uncertain)

    update_legend(slide)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
