from pathlib import Path
import shutil
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
REPO = ROOT.parent
TEMPLATE = REPO / "docs" / "template.pptx"
LOGO = ROOT / "generated-assets" / "txt-novigo-logo.png"
OUT = ROOT / "2026_GUBER_001 - Guber - NOVA Scenari infrastrutturali.pptx"

FONT = "Poppins"
W, H = 33.867, 19.05

C = {
    "navy": RGBColor(0x00, 0x44, 0x81),
    "blue": RGBColor(0x2A, 0x86, 0xCA),
    "sky": RGBColor(0x5B, 0xBE, 0xFE),
    "teal": RGBColor(0x2D, 0xCC, 0xCD),
    "teal2": RGBColor(0x7F, 0xE2, 0xE2),
    "pale": RGBColor(0xE9, 0xF7, 0xFA),
    "pale_blue": RGBColor(0xEC, 0xF5, 0xFC),
    "text": RGBColor(0x25, 0x2D, 0x36),
    "muted": RGBColor(0x67, 0x73, 0x7E),
    "line": RGBColor(0xC8, 0xDC, 0xE5),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "gray": RGBColor(0xF3, 0xF6, 0xF8),
    "green": RGBColor(0x47, 0x98, 0x79),
    "amber": RGBColor(0xD8, 0xA1, 0x43),
    "red": RGBColor(0xB5, 0x65, 0x65),
    "dark_teal": RGBColor(0x00, 0x7A, 0x82),
}


SCENARIOS = [
    {
        "num": "01", "name": "Infrastruttura attuale", "type": "Primario",
        "profile": "Continuità e cambiamento minimo",
        "arch": ["Guber", "Accesso SQL governato", "NOVA attuale"],
        "pros": ["Tempi più rapidi", "Rischio migrazione contenuto", "Alta reversibilità"],
        "attention": ["Ownership Guber limitata", "Accesso da segregare", "Dipendenza dall'ambiente corrente"],
        "gates": ["Viste / schema read-only", "Audit e capacità", "SLA operativo"],
        "oneoff": "€ 15k–38k", "monthly": "€ 0,7k–2,4k", "tco": "€ 40k–124k",
        "scores": [5, 5, 1, 4, 3, 4, 5],
    },
    {
        "num": "02", "name": "Primario su Azure Guber", "type": "Primario",
        "profile": "Ownership Guber e governance Azure",
        "arch": ["Utenti / sistemi", "Azure Guber", "DB NOVA primario"],
        "pros": ["Account e chiavi Guber", "Elasticità", "PaaS se compatibile"],
        "attention": ["Migrazione del primario", "Dipendenza dalla rete", "Run da assegnare"],
        "gates": ["Landing zone", "Compatibilità DB", "Networking privato"],
        "oneoff": "€ 49k–128k", "monthly": "€ 1,75k–5,3k", "tco": "€ 112k–319k",
        "scores": [2, 2, 5, 5, 5, 3, 3],
    },
    {
        "num": "03", "name": "Primario su AWS Guber", "type": "Primario",
        "profile": "Ownership Guber e target AWS Francoforte",
        "arch": ["Utenti / sistemi", "AWS eu-central-1", "DB NOVA primario"],
        "pros": ["Account e chiavi Guber", "Servizi gestiti", "Coerenza con S3"],
        "attention": ["Migrazione del primario", "Compatibilità RDS", "Responsabilità operative"],
        "gates": ["Standard AWS Guber", "VPC e connettività", "Motore / licensing"],
        "oneoff": "€ 49k–128k", "monthly": "€ 1,65k–5,2k", "tco": "€ 108k–315k",
        "scores": [2, 2, 5, 5, 5, 3, 3],
    },
    {
        "num": "04", "name": "Primario on-premise Guber", "type": "Primario",
        "profile": "Controllo infrastrutturale diretto",
        "arch": ["Utenti / sistemi", "Data center Guber", "DB NOVA primario"],
        "pros": ["Ownership Guber", "Integrazione con rete interna", "Controllo diretto"],
        "attention": ["Capacità e licensing", "HA / DR", "Carico operativo"],
        "gates": ["Infrastruttura disponibile", "Competenze", "Backup e secondo sito"],
        "oneoff": "€ 71k–208k", "monthly": "€ 1,4k–4,5k", "tco": "€ 121k–370k",
        "scores": [1, 1, 5, 5, 2, 1, 2],
    },
    {
        "num": "05", "name": "Replica on-premise Guber", "type": "Replica",
        "profile": "Autonomia sui dati senza spostare il transazionale",
        "arch": ["Primario attuale", "Replica asincrona", "On-premise Guber"],
        "pros": ["Migrazione contenuta", "Accesso locale", "Alta reversibilità"],
        "attention": ["Lag e riconciliazione", "Ownership divisa", "Compatibilità replica"],
        "gates": ["Motore / licenze", "Frequenza e SLA", "Viste e monitoring"],
        "oneoff": "€ 25k–64k", "monthly": "€ 0,9k–3,2k", "tco": "€ 57k–179k",
        "scores": [4, 4, 4, 5, 2, 3, 5],
    },
    {
        "num": "06", "name": "Replica su Azure Guber", "type": "Replica",
        "profile": "Accesso governato ai dati con servizi Azure",
        "arch": ["Primario attuale", "Replica asincrona", "Azure Guber"],
        "pros": ["Autonomia sui dati", "Identity e audit Azure", "Scalabilità"],
        "attention": ["Connettività cross-environment", "Lag ed egress", "Compatibilità replica"],
        "gates": ["Subscription e rete privata", "Replica supportata", "Retention"],
        "oneoff": "€ 25k–64k", "monthly": "€ 0,9k–3,2k", "tco": "€ 57k–179k",
        "scores": [4, 4, 4, 5, 4, 3, 5],
    },
]


def ensure_logo():
    if LOGO.exists():
        return
    with zipfile.ZipFile(TEMPLATE) as zf:
        with zf.open("ppt/media/image15.png") as src:
            LOGO.write_bytes(src.read())


def clear_slides(prs):
    ids = prs.slides._sldIdLst
    for slide_id in list(ids):
        prs.part.drop_rel(slide_id.rId)
        ids.remove(slide_id)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def shape(slide, x, y, w, h, fill="white", line="line", rounded=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = C[fill]
    s.line.color.rgb = C[line]
    s.line.width = Pt(0.8)
    return s


def text(slide, x, y, w, h, value, size=12, color="text", bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margins=0.04):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(margins)
    tf.margin_right = Cm(margins)
    tf.margin_top = Cm(margins)
    tf.margin_bottom = Cm(margins)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = C[color]
    return box


def bullets(slide, x, y, w, h, values, size=10, color="text", bullet_color=None, gap=5):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.03)
    tf.margin_bottom = Cm(0.03)
    for i, value in enumerate(values):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {value}"
        p.space_after = Pt(gap)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = C[color]
    return box


def line(slide, x1, y1, x2, y2, color="line", width=1.2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    c.line.color.rgb = C[color]
    c.line.width = Pt(width)
    return c


def header(slide, section, page):
    slide.shapes.add_picture(str(LOGO), Cm(0.65), Cm(0.35), width=Cm(2.55))
    text(slide, 3.45, 0.42, 20, 0.35, section.upper(), size=7.5, color="muted", bold=True)
    shape(slide, 0.67, 1.02, 1.35, 0.055, fill="teal", line="teal")


def title(slide, value, subtitle=None):
    text(slide, 1.0, 1.35, 31.4, 0.95, value, size=22, color="navy", bold=True)
    if subtitle:
        text(slide, 1.02, 2.28, 30.8, 0.55, subtitle, size=9.5, color="muted")


def card(slide, x, y, w, h, label, values, accent="teal", body_size=9):
    shape(slide, x, y, w, h, fill="white", line="line", rounded=True)
    shape(slide, x, y, w, 0.52, fill=accent, line=accent)
    text(slide, x + 0.2, y + 0.11, w - 0.4, 0.28, label, size=8.5, color="white", bold=True)
    if isinstance(values, list):
        bullets(slide, x + 0.25, y + 0.78, w - 0.5, h - 0.95, values, size=body_size, gap=4)
    else:
        text(slide, x + 0.25, y + 0.75, w - 0.5, h - 0.95, values, size=body_size)


def pill(slide, x, y, w, value, fill="pale", color="navy"):
    shape(slide, x, y, w, 0.55, fill=fill, line=fill, rounded=True)
    text(slide, x + 0.12, y + 0.13, w - 0.24, 0.25, value, size=8, color=color,
         bold=True, align=PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h, rows, widths=None, font=8.5, header_fill="navy",
              first_col_fill=None, heat=None):
    r_count, c_count = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(r_count, c_count, Cm(x), Cm(y), Cm(w), Cm(h)).table
    if widths:
        for idx, width in enumerate(widths):
            tbl.columns[idx].width = Cm(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Cm(0.08)
            cell.margin_right = Cm(0.08)
            cell.margin_top = Cm(0.04)
            cell.margin_bottom = Cm(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = C[header_fill]
                txt_color, bold = C["white"], True
            elif heat and c > 0:
                score = heat[r - 1][c - 1]
                palette = {1: "gray", 2: "pale_blue", 3: "pale", 4: "teal2", 5: "teal"}
                cell.fill.fore_color.rgb = C[palette[score]]
                txt_color, bold = C["navy"], score >= 4
            elif c == 0 and first_col_fill:
                cell.fill.fore_color.rgb = C[first_col_fill]
                txt_color, bold = C["navy"], True
            else:
                cell.fill.fore_color.rgb = C["white"] if r % 2 else C["gray"]
                txt_color, bold = C["text"], c == 0
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font)
                    run.font.bold = bold
                    run.font.color.rgb = txt_color
    return tbl


def cover(prs):
    s = blank(prs)
    s.shapes.add_picture(str(LOGO), Cm(0.8), Cm(0.55), width=Cm(2.8))
    shape(s, 0, 0, 0.22, H, fill="teal", line="teal")
    for x, y, d, col in [(26.1, 0.9, 5.8, "pale"), (28.0, 2.8, 3.6, "teal2"), (24.8, 5.0, 2.0, "sky")]:
        ring = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(x), Cm(y), Cm(d), Cm(d))
        ring.fill.background()
        ring.line.color.rgb = C[col]
        ring.line.width = Pt(2)
    pill(s, 1.2, 6.4, 5.2, "DECISION PACK", fill="navy", color="white")
    text(s, 1.2, 7.35, 25.5, 1.65, "NOVA — alternative infrastrutturali\ne modelli operativi", size=30, color="navy", bold=True)
    text(s, 1.25, 10.2, 22, 0.65, "Confronto preliminare per Guber", size=14, color="blue", bold=True)
    text(s, 1.25, 16.9, 10, 0.4, "31 luglio 2026", size=9, color="muted")
    text(s, 24.0, 16.9, 8.5, 0.4, "Range di pianificazione", size=8, color="muted", align=PP_ALIGN.RIGHT)


def context(prs, page):
    s = blank(prs); header(s, "Contesto / Esigenza / Obiettivi", page)
    title(s, "Il perimetro combina accesso ai dati, controllo e continuità",
          "La replica read-only presso Guber è considerata accettabile; la scelta della collocazione resta aperta.")
    items = [
        ("01", "Contesto", "Guber richiede accesso SQL ai dati NOVA per verifiche ed estrazioni."),
        ("02", "Esigenza", "Rendere i dati accessibili senza compromettere operatività, sicurezza e responsabilità."),
        ("03", "Obiettivo", "Confrontare sei alternative di collocazione e gestione, senza anticipare la scelta."),
    ]
    for i, (n, label, body) in enumerate(items):
        x = 1.0 + i * 10.65
        shape(s, x, 4.0, 9.6, 8.9, fill="white", line="line", rounded=True)
        shape(s, x + 0.45, 4.55, 1.35, 1.35, fill="teal" if i != 1 else "blue", line="teal", rounded=True)
        text(s, x + 0.45, 4.9, 1.35, 0.35, n, size=11, color="white", bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.45, 6.5, 8.5, 0.6, label, size=16, color="navy", bold=True)
        text(s, x + 0.45, 7.5, 8.35, 2.7, body, size=12, color="text")
    shape(s, 1.0, 14.2, 31.0, 1.7, fill="pale", line="pale", rounded=True)
    text(s, 1.45, 14.67, 30.1, 0.65,
         "Punto fermo: l'accesso ai dati è confermato; la necessità di migrare il database primario non è ancora dimostrata.",
         size=12, color="dark_teal", bold=True, align=PP_ALIGN.CENTER)


def asis(prs, page):
    s = blank(prs); header(s, "AS IS", page)
    title(s, "Oggi NOVA concentra dati e operatività sull'ambiente corrente",
          "Il requisito SQL è chiaro; dimensionamento e vincoli infrastrutturali devono ancora essere validati.")
    boxes = [
        (1.2, "Sorgenti", "ODT · AV · LDT"),
        (8.1, "NOVA", "History e consolidamento"),
        (15.0, "Database primario", "Operatività applicativa"),
        (23.2, "Guber", "Verifiche ed estrazioni SQL"),
    ]
    for i, (x, label, body) in enumerate(boxes):
        shape(s, x, 5.0, 5.8 if i != 2 else 7.0, 2.5, fill="white", line="teal", rounded=True)
        text(s, x + 0.25, 5.5, 5.3 if i != 2 else 6.5, 0.4, label, size=12, color="navy", bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.25, 6.25, 5.3 if i != 2 else 6.5, 0.45, body, size=9, color="muted", align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            end = x + (5.8 if i != 2 else 7.0)
            next_x = boxes[i + 1][0]
            line(s, end, 6.25, next_x, 6.25, "teal", 1.8)
    card(s, 1.2, 9.4, 9.5, 4.0, "Fatto confermato", [
        "Accesso SQL richiesto",
        "Replica read-only accettata",
        "AM Novigo end-to-end disponibile",
    ], "teal", 10)
    card(s, 11.45, 9.4, 9.5, 4.0, "Da proteggere", [
        "Disponibilità del transazionale",
        "History e riconciliazione",
        "Accessi, audit e segregazione",
    ], "blue", 10)
    card(s, 21.7, 9.4, 9.5, 4.0, "Dati mancanti", [
        "Motore, sizing e licensing",
        "SLA, RPO/RTO e rete",
        "Standard cloud e capacità Guber",
    ], "navy", 10)


def scenario_map(prs, page):
    s = blank(prs); header(s, "Alternative", page)
    title(s, "Sei alternative separano due scelte: collocare il primario o la replica",
          "La mappa evita di confondere il requisito di accesso ai dati con una migrazione completa della piattaforma.")
    shape(s, 13.0, 3.3, 7.9, 1.2, fill="navy", line="navy", rounded=True)
    text(s, 13.2, 3.68, 7.5, 0.35, "COME SODDISFARE IL REQUISITO?", size=10, color="white", bold=True, align=PP_ALIGN.CENTER)
    line(s, 17.0, 4.5, 17.0, 5.3, "teal", 2)
    line(s, 8.5, 5.3, 25.5, 5.3, "teal", 2)
    for x, label in [(8.5, "COLLOCAZIONE DEL PRIMARIO"), (25.5, "PRIMARIO ATTUALE + REPLICA")]:
        line(s, x, 5.3, x, 6.1, "teal", 2)
        shape(s, x - 3.4, 6.1, 6.8, 1.05, fill="pale", line="teal", rounded=True)
        text(s, x - 3.1, 6.42, 6.2, 0.32, label, size=9, color="navy", bold=True, align=PP_ALIGN.CENTER)
    primary = SCENARIOS[:4]
    for i, sc in enumerate(primary):
        x = 0.9 + i * 7.85
        shape(s, x, 8.2, 7.0, 4.6, fill="white", line="line", rounded=True)
        pill(s, x + 0.35, 8.55, 1.1, sc["num"], fill="teal", color="white")
        text(s, x + 0.35, 9.55, 6.3, 1.0, sc["name"], size=11, color="navy", bold=True)
        text(s, x + 0.35, 11.0, 6.3, 0.9, sc["profile"], size=8.5, color="muted")
    for i, sc in enumerate(SCENARIOS[4:]):
        x = 17.4 + i * 7.85
        shape(s, x, 13.45, 7.0, 3.3, fill="white", line="teal", rounded=True)
        pill(s, x + 0.35, 13.8, 1.1, sc["num"], fill="blue", color="white")
        text(s, x + 0.35, 14.75, 6.3, 0.8, sc["name"], size=10.5, color="navy", bold=True)
        text(s, x + 0.35, 15.65, 6.3, 0.45, "Replica read-only presso Guber", size=8, color="muted")


def scenario_slide(prs, page, sc):
    s = blank(prs); header(s, f"Scenario {sc['num']}", page)
    title(s, f"{sc['name']}: {sc['profile'].lower()}",
          "Profilo preliminare; fattibilità e costo dipendono dai gate tecnici indicati.")
    pill(s, 28.2, 1.4, 3.7, sc["type"].upper(), fill="pale_blue")
    # Architecture strip
    for i, label in enumerate(sc["arch"]):
        x = 1.1 + i * 9.8
        fill = "pale" if i == 1 else "white"
        shape(s, x, 3.6, 7.7, 1.55, fill=fill, line="teal", rounded=True)
        text(s, x + 0.25, 4.05, 7.2, 0.45, label, size=10, color="navy", bold=True, align=PP_ALIGN.CENTER)
        if i < 2:
            line(s, x + 7.7, 4.38, x + 9.45, 4.38, "teal", 1.8)
    card(s, 1.1, 6.2, 9.55, 5.0, "PUNTI DI FORZA", sc["pros"], "teal", 10)
    card(s, 11.25, 6.2, 9.55, 5.0, "ATTENZIONI", sc["attention"], "blue", 10)
    card(s, 21.4, 6.2, 9.55, 5.0, "GATE DA VALIDARE", sc["gates"], "navy", 10)
    shape(s, 1.1, 12.25, 29.85, 2.45, fill="pale", line="pale", rounded=True)
    text(s, 1.6, 12.68, 8.2, 0.35, "UNA TANTUM", size=8, color="muted", bold=True, align=PP_ALIGN.CENTER)
    text(s, 1.6, 13.25, 8.2, 0.55, sc["oneoff"], size=17, color="navy", bold=True, align=PP_ALIGN.CENTER)
    text(s, 11.3, 12.68, 8.2, 0.35, "PIATTAFORMA / MESE", size=8, color="muted", bold=True, align=PP_ALIGN.CENTER)
    text(s, 11.3, 13.25, 8.2, 0.55, sc["monthly"], size=17, color="navy", bold=True, align=PP_ALIGN.CENTER)
    text(s, 21.0, 12.68, 8.2, 0.35, "TCO PIATTAFORMA · 3 ANNI", size=8, color="muted", bold=True, align=PP_ALIGN.CENTER)
    text(s, 21.0, 13.25, 8.2, 0.55, sc["tco"], size=17, color="dark_teal", bold=True, align=PP_ALIGN.CENTER)
    text(s, 1.25, 15.35, 29.5, 0.7,
         "Range di pianificazione · AM Novigo, team sostitutivo, licenze non note e DR geografico esclusi.",
         size=8.5, color="muted", align=PP_ALIGN.CENTER)


def operating_model(prs, page):
    s = blank(prs); header(s, "Modello operativo", page)
    title(s, "Il run operativo resta una scelta trasversale a tutti gli scenari",
          "Ownership di account e chiavi non implica che Guber debba eseguire direttamente tutte le attività operative.")
    card(s, 1.2, 4.0, 14.8, 8.8, "AM NOVIGO END-TO-END", [
        "Un presidio operativo coordinato",
        "Applicazione, database e infrastruttura",
        "Monitoring, backup, patching e sicurezza",
        "Incident, capacity e release management",
        "Accessi delegati e tracciati",
    ], "teal", 11)
    card(s, 17.0, 4.0, 14.8, 8.8, "GUBER / TERZA PARTE", [
        "Maggiore autonomia di sourcing",
        "Competenze e copertura da garantire",
        "Handover e documentazione operativa",
        "Rischio di responsabilità frammentate",
        "Coordinamento multi-fornitore",
    ], "blue", 11)
    shape(s, 6.2, 14.1, 21.4, 1.4, fill="pale", line="pale", rounded=True)
    text(s, 6.6, 14.52, 20.6, 0.55,
         "In entrambi i modelli servono RACI, SLA, key ownership, audit ed exit plan.",
         size=11, color="dark_teal", bold=True, align=PP_ALIGN.CENTER)


def economics(prs, page):
    s = blank(prs); header(s, "Economics", page)
    title(s, "I range delimitano il confronto, ma non sostituiscono il sizing",
          "TCO piattaforma = una tantum + 36 mesi; Application Maintenance valorizzata separatamente.")
    rows = [["Scenario", "Una tantum", "Piattaforma / mese", "TCO 3 anni"]]
    for sc in SCENARIOS:
        rows.append([sc["name"], sc["oneoff"], sc["monthly"], sc["tco"]])
    add_table(s, 1.0, 3.5, 31.1, 9.4, rows, widths=[12.8, 5.7, 6.2, 6.4], font=9.2, first_col_fill="pale")
    notes = [
        "AWS verificato su Francoforte, eu-central-1",
        "AM e team sostitutivo esclusi",
        "Licenze non note e DR geografico esclusi",
        "Confidenza bassa fino al sizing reale",
    ]
    for i, n in enumerate(notes):
        pill(s, 1.0 + i * 7.8, 14.0, 7.2, n, fill="pale_blue")


def matrix(prs, page):
    s = blank(prs); header(s, "Confronto", page)
    title(s, "La matrice evidenzia trade-off differenti, non un vincitore",
          "Lettura qualitativa preliminare: da rivedere dopo la verifica dei gate tecnici e degli standard Guber.")
    criteria = ["Rapidità", "Rischio\ncontenuto", "Ownership\nGuber", "Accesso\nSQL", "Elasticità", "Semplicità\noperativa", "Reversibilità"]
    rows = [["Scenario"] + criteria]
    heat = []
    labels = {1: "Bassa", 2: "Medio-\nbassa", 3: "Media", 4: "Medio-\nalta", 5: "Alta"}
    for sc in SCENARIOS:
        rows.append([sc["name"]] + [labels[v] for v in sc["scores"]])
        heat.append(sc["scores"])
    add_table(s, 0.8, 3.65, 32.2, 10.0, rows,
              widths=[10.8] + [3.05] * 7, font=8.2, first_col_fill="pale", heat=heat)
    text(s, 1.0, 14.55, 31.0, 0.7,
         "Valutazione preliminare: standard Guber, compatibilità del motore e capacità on-premise possono modificare il posizionamento.",
         size=9, color="muted", align=PP_ALIGN.CENTER)


def synthesis(prs, page):
    s = blank(prs); header(s, "Sintesi", page)
    title(s, "Le alternative restano aperte; cambiano le informazioni necessarie",
          "Il deck presenta tre famiglie di scelta senza attribuire oggi una preferenza.")
    groups = [
        ("CONTINUITÀ", "Scenario 1", "Mantenere il primario sull'ambiente corrente", "teal"),
        ("TRASFERIMENTO DEL PRIMARIO", "Scenari 2–4", "Azure, AWS oppure on-premise Guber", "blue"),
        ("ACCESSO DATI SENZA MIGRAZIONE", "Scenari 5–6", "Replica on-premise oppure Azure Guber", "navy"),
    ]
    for i, (label, ref, body, col) in enumerate(groups):
        x = 1.0 + i * 10.65
        shape(s, x, 3.8, 9.6, 6.4, fill="white", line="line", rounded=True)
        shape(s, x, 3.8, 9.6, 0.65, fill=col, line=col)
        text(s, x + 0.35, 4.8, 8.9, 0.45, label, size=9, color="muted", bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.35, 5.75, 8.9, 0.55, ref, size=16, color="navy", bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.7, 7.0, 8.2, 1.4, body, size=11, color="text", align=PP_ALIGN.CENTER)
    shape(s, 1.0, 11.5, 31.0, 3.6, fill="pale", line="pale", rounded=True)
    text(s, 1.5, 12.0, 30, 0.4, "INFORMAZIONI CHE RIDUCONO L'INCERTEZZA", size=9, color="dark_teal", bold=True, align=PP_ALIGN.CENTER)
    items = ["Motore e sizing", "Standard cloud", "Capacità on-premise", "SLA / RPO / RTO", "Connettività", "Licensing", "Modello operativo"]
    for i, item in enumerate(items):
        pill(s, 1.55 + i * 4.25, 13.05, 3.8, item, fill="white")


def appendix_title(prs, page):
    s = blank(prs); header(s, "Appendice", page)
    shape(s, 0, 1.15, W, 17.9, fill="navy", line="navy")
    text(s, 1.2, 6.3, 23, 0.6, "APPENDICE", size=11, color="teal2", bold=True)
    text(s, 1.2, 7.3, 26, 1.3, "Assunzioni, economics e\napprofondimenti tecnici", size=28, color="white", bold=True)
    text(s, 1.2, 10.7, 23, 0.6, "Materiale di supporto — range preliminari", size=11, color="teal2")


def evidence(prs, page):
    s = blank(prs); header(s, "Appendice | Grounding", page)
    title(s, "Fatti, deduzioni e dati mancanti restano separati")
    card(s, 1.0, 3.3, 9.8, 10.3, "FATTI CONFERMATI", [
        "Accesso SQL richiesto",
        "Replica read-only accettata",
        "Azure e on-premise in agenda",
        "AWS aggiunto su conferma utente",
        "AM Novigo end-to-end",
    ], "teal", 10)
    card(s, 12.0, 3.3, 9.8, 10.3, "DEDUZIONI DICHIARATE", [
        "Migrare il primario non è l'unico modo",
        "La replica riduce il rischio sul transazionale",
        "DB remoto rende la rete parte dell'architettura",
        "Key ownership è una scelta di governance",
    ], "blue", 10)
    card(s, 23.0, 3.3, 9.8, 10.3, "DATI MANCANTI", [
        "Motore, versione e licensing",
        "Sizing, crescita e IOPS",
        "SLA, RPO / RTO e downtime",
        "Landing zone e connettività",
        "Capacità on-premise e contratti",
    ], "navy", 10)


def assumptions(prs, page):
    s = blank(prs); header(s, "Appendice | Economics", page)
    title(s, "Il modello economico usa un carico di riferimento esplicito",
          "Le assunzioni rendono comparabili gli scenari, ma non sostituiscono i dati NOVA.")
    rows = [
        ["Voce", "Assunzione"],
        ["Produzione", "4–8 vCPU · 16–32 GB RAM · 24x7 · HA"],
        ["Dati DB", "500 GB iniziali · crescita 20% annuo"],
        ["Non produzione", "1 ambiente al 50% della produzione"],
        ["Backup", "30 giorni · restore periodicamente testato"],
        ["Documentale", "1 TB iniziale"],
        ["Accesso SQL", "20 utenti autorizzati · max 5 concorrenti"],
        ["TCO", "Una tantum + 36 mesi di piattaforma"],
    ]
    add_table(s, 1.1, 3.5, 20.5, 10.5, rows, widths=[6.0, 14.5], font=9.5, first_col_fill="pale")
    card(s, 22.6, 3.5, 9.4, 10.5, "NON INCLUSO", [
        "AM Novigo",
        "Team sostitutivo",
        "Licenze non note",
        "DR geografico",
        "Costi interni Guber",
        "IVA e sconti contrattuali",
    ], "blue", 10)


def cost_detail(prs, page, sc):
    s = blank(prs); header(s, f"Appendice | Scenario {sc['num']}", page)
    title(s, f"{sc['name']} — struttura del range economico",
          "Componenti indicative da ricalcolare dopo assessment e sizing.")
    oneoff_items = {
        "01": ["Assessment e hardening", "Accesso SQL / viste", "Audit e test", "Documentazione"],
        "02": ["Assessment PaaS / IaaS", "Landing zone e rete", "Migrazione e test", "Cutover e rollback"],
        "03": ["Assessment RDS / EC2", "Landing zone e VPC", "Migrazione e test", "Cutover e rollback"],
        "04": ["Assessment e capacity", "Provisioning e licensing", "Migrazione e test", "HA, backup e cutover"],
        "05": ["Assessment replica", "Provisioning on-premise", "Setup e riconciliazione", "Accessi e monitoring"],
        "06": ["Assessment replica", "Subscription e rete", "Setup e riconciliazione", "Accessi e monitoring"],
    }[sc["num"]]
    recurring = {
        "01": ["Capacità incrementale", "Backup e logging", "Monitoring", "Networking"],
        "02": ["Compute / database", "Storage e backup", "Logging e security", "Networking"],
        "03": ["RDS / EC2", "Storage e backup", "CloudWatch / security", "Networking"],
        "04": ["Capacità e manutenzione", "Backup", "Monitoring", "Energia / data center"],
        "05": ["Compute / database replica", "Storage e backup", "Monitoring del lag", "Connettività"],
        "06": ["Database replica", "Storage e backup", "Monitor e audit", "Rete / egress"],
    }[sc["num"]]
    card(s, 1.1, 3.5, 9.8, 8.7, "UNA TANTUM", oneoff_items, "teal", 10)
    card(s, 11.8, 3.5, 9.8, 8.7, "RICORRENTE", recurring, "blue", 10)
    card(s, 22.5, 3.5, 9.8, 8.7, "SENSIBILITÀ", [
        "Motore e compatibilità",
        "HA / DR e SLA",
        "Volumi, crescita e IOPS",
        "Licenze e contratti",
    ], "navy", 10)
    shape(s, 1.1, 13.25, 31.2, 1.8, fill="pale", line="pale", rounded=True)
    text(s, 1.55, 13.65, 9.6, 0.35, f"Una tantum  {sc['oneoff']}", size=11, color="navy", bold=True, align=PP_ALIGN.CENTER)
    text(s, 11.85, 13.65, 9.6, 0.35, f"Mensile  {sc['monthly']}", size=11, color="navy", bold=True, align=PP_ALIGN.CENTER)
    text(s, 22.15, 13.65, 9.6, 0.35, f"TCO 3 anni  {sc['tco']}", size=11, color="dark_teal", bold=True, align=PP_ALIGN.CENTER)


def security(prs, page):
    s = blank(prs); header(s, "Appendice | Sicurezza", page)
    title(s, "Account, chiavi e run possono avere ownership differenti")
    labels = [
        ("ACCOUNT E INFRASTRUTTURA", "Guber", "navy"),
        ("CHIAVI E POLICY", "Guber / doppio controllo", "blue"),
        ("RUN OPERATIVO", "Novigo oppure Guber / terzo", "teal"),
        ("AUDIT E GOVERNANCE", "Guber", "dark_teal"),
    ]
    for i, (label, owner, col) in enumerate(labels):
        x = 1.2 + i * 7.9
        shape(s, x, 4.3, 7.0, 5.4, fill="white", line=col, rounded=True)
        shape(s, x, 4.3, 7.0, 0.65, fill=col, line=col)
        text(s, x + 0.35, 5.45, 6.3, 1.0, label, size=9, color="muted", bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.45, 7.1, 6.1, 1.1, owner, size=14, color="navy", bold=True, align=PP_ALIGN.CENTER)
    card(s, 1.2, 11.2, 30.7, 3.3, "CONTROLLO MINIMO COMUNE", [
        "TLS e cifratura a riposo · accesso privato · MFA e federation · logging e SIEM · backup protetti · rotazione e revoca delle chiavi",
    ], "teal", 10)


def storage(prs, page):
    s = blank(prs); header(s, "Appendice | Storage", page)
    title(s, "La tecnologia di storage dipende dal requisito, non dal nome del provider",
          "Prima di scegliere S3 va chiarito se il requisito è AWS-specifico o genericamente object storage.")
    options = [
        ("AWS S3", ["Coerenza con AWS", "Lifecycle e versioning", "Valutare egress e KMS"], "teal"),
        ("Azure Storage", ["Coerenza con Azure", "Identity e policy", "Valutare tier e rete"], "blue"),
        ("Object storage on-premise", ["Controllo locale", "Capacity e operations", "Durabilità da progettare"], "navy"),
    ]
    for i, (name, vals, col) in enumerate(options):
        card(s, 1.1 + i * 10.5, 4.1, 9.4, 7.4, name.upper(), vals, col, 10)
    shape(s, 1.1, 12.7, 30.4, 2.0, fill="pale", line="pale", rounded=True)
    text(s, 1.5, 13.1, 29.6, 0.85,
         "Criteri comuni: volume e crescita · access pattern · retention · immutabilità · cifratura · data residency · costo di trasferimento.",
         size=10.5, color="dark_teal", bold=True, align=PP_ALIGN.CENTER)


def orchestrator(prs, page):
    s = blank(prs); header(s, "Appendice | Orchestratore", page)
    title(s, "Separare database e orchestratore rende la rete parte dell'architettura")
    paths = [
        ("DB-ONLY", "Il database migra o replica; applicazione e orchestratore restano dove sono.", ["Più rapido", "Dipendenza WAN"], "teal"),
        ("PHASED", "Replica o DB come prima fase; componenti spostati dopo verifica.", ["Rischio progressivo", "Transitorio da governare"], "blue"),
        ("FULL-STACK", "Applicazione, database e orchestratore migrano insieme.", ["Coerenza target", "Scope e cutover maggiori"], "navy"),
    ]
    for i, (name, body, tags, col) in enumerate(paths):
        x = 1.0 + i * 10.65
        shape(s, x, 3.8, 9.6, 8.1, fill="white", line=col, rounded=True)
        shape(s, x, 3.8, 9.6, 0.65, fill=col, line=col)
        text(s, x + 0.4, 4.9, 8.8, 0.45, name, size=13, color="navy", bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.6, 6.1, 8.4, 2.0, body, size=10, color="text", align=PP_ALIGN.CENTER)
        pill(s, x + 0.7, 9.0, 3.8, tags[0], fill="pale")
        pill(s, x + 5.0, 9.0, 3.8, tags[1], fill="pale_blue")
    text(s, 1.0, 13.2, 31.0, 0.8,
         "Un'architettura DB-only può essere una fase controllata; non deve diventare implicitamente un target fragile.",
         size=11, color="dark_teal", bold=True, align=PP_ALIGN.CENTER)


def paas(prs, page):
    s = blank(prs); header(s, "Appendice | Deployment DB", page)
    title(s, "PaaS o IaaS è un gate di compatibilità, non un dettaglio implementativo")
    card(s, 1.2, 4.0, 14.7, 8.8, "SERVIZIO GESTITO · PaaS / RDS", [
        "Patching e backup maggiormente gestiti",
        "HA e monitoring standardizzati",
        "Minor carico operativo",
        "Compatibilità del motore da verificare",
        "Vincoli e feature specifiche del servizio",
    ], "teal", 10.5)
    card(s, 17.0, 4.0, 14.7, 8.8, "VM / EC2 · SELF-MANAGED", [
        "Maggiore compatibilità e controllo",
        "Responsabilità DBA e sistema operativo",
        "Backup, HA e patching da costruire",
        "TCO operativo più elevato",
        "Licensing da verificare",
    ], "blue", 10.5)
    pill(s, 10.5, 14.1, 12.8, "PRIMO GATE: MOTORE · VERSIONE · FEATURE · LICENZE", fill="navy", color="white")


def raci(prs, page):
    s = blank(prs); header(s, "Appendice | RACI", page)
    title(s, "Il modello di run deve coprire l'intero servizio, non soltanto l'applicazione")
    rows = [
        ["Ambito", "AM Novigo end-to-end", "Guber / terza parte"],
        ["Applicazione e release", "Novigo esegue", "Owner designato esegue"],
        ["Database e performance", "Novigo esegue", "DBA Guber / terzo"],
        ["Infrastruttura / cloud", "Novigo opera con delega", "Guber / terzo"],
        ["Backup e restore", "Novigo esegue e testa", "Guber / terzo"],
        ["Security operations", "Novigo opera; Guber governa", "Guber / terzo"],
        ["Account e key ownership", "Guber governa", "Guber governa"],
        ["Incident e problem", "Presidio coordinato Novigo", "Coordinamento multi-owner"],
        ["Exit / handover", "Obbligo contrattuale", "Processo interno / sourcing"],
    ]
    add_table(s, 1.0, 3.6, 31.2, 10.7, rows, widths=[8.8, 11.2, 11.2], font=9, first_col_fill="pale")


def indicative_path(prs, page):
    s = blank(prs); header(s, "Appendice | Approfondimento", page)
    title(s, "Un eventuale approfondimento può ridurre l'incertezza in quattro passaggi",
          "Sequenza indicativa, non obiettivo immediato né impegno di progetto.")
    phases = [
        ("01", "Assessment", "1–2 settimane", "Motore, sizing, licenze e compatibilità"),
        ("02", "Economics", "1 settimana", "TCO, contratti e modello AM"),
        ("03", "Prova tecnica", "2–4 settimane", "Replica o migrazione pilota"),
        ("04", "Piano di transizione", "Da definire", "Cutover, rollback, RACI e SLA"),
    ]
    for i, (n, name, duration, output) in enumerate(phases):
        x = 1.1 + i * 7.85
        shape(s, x, 4.2, 7.0, 7.4, fill="white", line="line", rounded=True)
        pill(s, x + 0.35, 4.6, 1.2, n, fill="teal", color="white")
        text(s, x + 0.35, 5.75, 6.3, 0.5, name, size=13, color="navy", bold=True)
        text(s, x + 0.35, 6.65, 6.3, 0.4, duration, size=9, color="blue", bold=True)
        text(s, x + 0.35, 8.05, 6.2, 1.7, output, size=9.5, color="text")
        if i < 3:
            line(s, x + 7.0, 7.9, x + 7.7, 7.9, "teal", 1.8)
    shape(s, 1.1, 13.0, 30.6, 1.8, fill="pale", line="pale", rounded=True)
    text(s, 1.5, 13.42, 29.8, 0.7,
         "L'attivazione del percorso dipende dalla successiva scelta di approfondire una o più alternative.",
         size=10.5, color="dark_teal", bold=True, align=PP_ALIGN.CENTER)


def build():
    ensure_logo()
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)
    prs.slide_width = Cm(W)
    prs.slide_height = Cm(H)

    cover(prs)
    context(prs, 2)
    asis(prs, 3)
    scenario_map(prs, 4)
    page = 5
    for sc in SCENARIOS:
        scenario_slide(prs, page, sc)
        page += 1
    operating_model(prs, 11)
    economics(prs, 12)
    matrix(prs, 13)
    synthesis(prs, 14)

    appendix_title(prs, 15)
    evidence(prs, 16)
    assumptions(prs, 17)
    page = 18
    for sc in SCENARIOS:
        cost_detail(prs, page, sc)
        page += 1
    security(prs, 24)
    storage(prs, 25)
    orchestrator(prs, 26)
    paas(prs, 27)
    raci(prs, 28)
    indicative_path(prs, 29)

    props = prs.core_properties
    props.title = "NOVA — alternative infrastrutturali e modelli operativi"
    props.subject = "Decision pack preliminare per Guber"
    props.author = "TXT e-solutions / Novigo"
    props.comments = "Range di pianificazione; nessuna raccomandazione strategica."
    prs.save(OUT)
    print(OUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
