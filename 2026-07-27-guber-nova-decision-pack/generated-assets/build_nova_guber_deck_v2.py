from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

from build_visual_prototype_v2 import (
    C, FONT, LOGO, STAR, WAVE, box, bullets, connector, header, slide, textbox, title
)
from build_nova_guber_deck import SCENARIOS


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "2026_GUBER_001 - Guber - NOVA Scenari infrastrutturali v2.pptx"


SCENARIO_MESSAGES = {
    "01": "La continuità riduce il rischio di transizione, ma limita il controllo infrastrutturale di Guber",
    "02": "Azure trasferisce il controllo a Guber e richiede una migrazione governata del primario",
    "03": "AWS trasferisce il controllo a Guber e valorizza i servizi della regione Francoforte",
    "04": "L'on-premise massimizza il controllo diretto e concentra su Guber capacità e continuità operativa",
    "05": "La replica on-premise abilita l'accesso locale senza spostare il database transazionale",
    "06": "La replica Azure abilita accesso e governance cloud senza migrare il database transazionale",
}

SCENARIO_DESCRIPTIONS = {
    "01": "Il database primario resta nell'ambiente corrente; l'accesso SQL viene segregato e tracciato.",
    "02": "Il database primario migra nella subscription Azure di Guber, su servizio gestito o VM compatibile.",
    "03": "Il database primario migra nell'account AWS di Guber, in regione Europe (Frankfurt).",
    "04": "Il database primario migra nell'infrastruttura Guber, con HA, backup e operations da dimensionare.",
    "05": "Il primario resta nell'ambiente corrente e alimenta una replica read-only nell'infrastruttura Guber.",
    "06": "Il primario resta nell'ambiente corrente e alimenta una replica read-only nella subscription Azure Guber.",
}


def star(s):
    s.shapes.add_picture(str(STAR), Cm(31.9), Cm(17.3), width=Cm(0.65))


def add_runs(box_shape, segments, size=10, align=PP_ALIGN.LEFT):
    tf = box_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    for value, color, bold in segments:
        r = p.add_run()
        r.text = value
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = C[color]


def clean_cover(prs):
    s = slide(prs)
    s.shapes.add_picture(str(LOGO), Cm(0.65), Cm(0.55), width=Cm(3.2))
    s.shapes.add_picture(str(WAVE), Cm(22.7), Cm(0), height=Cm(19.05))
    textbox(s, 0.75, 12.8, 21.0, 0.7, "NOVA", 19, "teal", True)
    textbox(s, 0.75, 13.55, 21.0, 1.25, "Alternative infrastrutturali\ne modelli operativi", 22, "text", True)
    textbox(s, 0.75, 15.2, 16.0, 0.4, "Guber", 10, "gray")
    textbox(s, 0.75, 17.5, 8.0, 0.35, "31 luglio 2026", 8, "gray")


def context_slide(prs):
    s = slide(prs); header(s, "NOVA", "Contesto, esigenza e obiettivi", 2)
    xline = 5.2
    connector(s, xline, 3.0, xline, 15.6, "teal", 1.1)
    rows = [
        ("Contesto", "Guber richiede un accesso SQL governato ai dati NOVA",
         "L'esigenza riguarda verifiche ed estrazioni, mantenendo continuità operativa e tracciabilità degli accessi."),
        ("Esigenza", "Separare l'autonomia sui dati dal rischio sul transazionale",
         "La collocazione deve bilanciare controllo Guber, sicurezza, responsabilità operative e impatto della transizione."),
        ("Obiettivi", "Confrontare sei alternative su criteri omogenei",
         "Il confronto considera collocazione, accesso SQL, modello operativo, rischio, reversibilità ed economics."),
    ]
    for i, (label, lead, body) in enumerate(rows):
        y = 3.0 + i * 4.35
        textbox(s, 1.1, y + 0.3, 3.6, 0.4, label, 10, "dark_teal", True, PP_ALIGN.RIGHT)
        textbox(s, 6.2, y, 24.7, 0.55, lead, 12, "text", True)
        textbox(s, 6.2, y + 0.85, 24.7, 1.25, body, 9.5, "text")
    star(s)


def architecture_node(s, x, y, w, label, sub="", fill="white"):
    box(s, x, y, w, 1.45, fill, "line", True)
    textbox(s, x + 0.2, y + 0.3, w - 0.4, 0.35, label, 9.5, "navy", True, PP_ALIGN.CENTER)
    if sub:
        textbox(s, x + 0.2, y + 0.85, w - 0.4, 0.25, sub, 7.3, "gray", False, PP_ALIGN.CENTER)


def as_is_slide(prs):
    s = slide(prs); header(s, "NOVA", "Situazione attuale e requisito dati", 3)
    title(s, "Il requisito SQL si innesta su una piattaforma oggi concentrata nell'ambiente corrente")
    textbox(s, 0.8, 2.4, 31.0, 0.45,
            "NOVA gestisce history e consolidamento; l'accesso Guber deve essere separato dal workload applicativo.", 9, "gray")
    box(s, 0.9, 3.5, 16.0, 10.7, "white", "teal", True)
    textbox(s, 1.4, 3.9, 15.0, 0.4, "FLUSSO OPERATIVO", 8, "dark_teal", True, PP_ALIGN.CENTER)
    nodes = [
        ("Sorgenti", "ODT · AV · Loan Data Tape"),
        ("NOVA", "history · working row · consolidamento"),
        ("Database primario", "servizio applicativo"),
    ]
    for i, (label, sub) in enumerate(nodes):
        y = 5.0 + i * 2.45
        architecture_node(s, 4.1, y, 9.6, label, sub, "pale_blue" if i == 1 else "white")
        if i < 2:
            connector(s, 8.9, y + 1.45, 8.9, y + 2.35)
    textbox(s, 18.3, 3.8, 12.0, 0.4, "REQUISITO GUBER", 9, "dark_teal", True)
    bullets(s, 18.3, 4.5, 13.0, [
        "Interrogazioni SQL per verifiche ed estrazioni",
        "Profili read-only come configurazione ordinaria",
        "Accessi, query e operazioni amministrative tracciati",
    ], 9.5)
    textbox(s, 18.3, 9.2, 12.0, 0.4, "IMPLICAZIONE ARCHITETTURALE", 9, "blue", True)
    bullets(s, 18.3, 9.9, 13.0, [
        "Il requisito può essere soddisfatto sul primario o su una replica",
        "Una replica evita query dirette sul transazionale",
        "Frequenza, lag e oggetti esposti diventano requisiti di servizio",
    ], 9.5)
    star(s)


def alternatives_slide(prs):
    s = slide(prs); header(s, "NOVA", "Alternative di collocazione", 4)
    title(s, "Le sei alternative distinguono la collocazione del primario dalla replica presso Guber")
    textbox(s, 1.0, 2.45, 31.0, 0.45,
            "Quattro opzioni intervengono sul database primario; due mantengono il primario e aggiungono una replica read-only.", 9, "gray")
    # Primary family
    textbox(s, 1.0, 3.45, 18.9, 0.4, "COLLOCAZIONE DEL DATABASE PRIMARIO", 9, "dark_teal", True)
    connector(s, 1.0, 4.0, 20.4, 4.0, "teal", 1.3)
    for i, sc in enumerate(SCENARIOS[:4]):
        x = 1.0 + i * 7.75
        box(s, x, 4.75, 6.8, 5.15, "white", "line", True)
        textbox(s, x + 0.35, 5.15, 0.8, 0.35, sc["num"], 9, "dark_teal", True)
        textbox(s, x + 0.35, 6.0, 6.1, 0.85, sc["name"], 11, "navy", True)
        textbox(s, x + 0.35, 7.35, 6.1, 1.2, sc["profile"], 8.5, "gray")
    # Replica family
    textbox(s, 1.0, 11.1, 18.9, 0.4, "PRIMARIO ATTUALE + REPLICA READ-ONLY", 9, "blue", True)
    connector(s, 1.0, 11.65, 20.4, 11.65, "blue", 1.3)
    for i, sc in enumerate(SCENARIOS[4:]):
        x = 8.9 + i * 8.5
        box(s, x, 12.4, 7.4, 3.4, "white", "teal", True)
        textbox(s, x + 0.35, 12.75, 0.8, 0.35, sc["num"], 9, "blue", True)
        textbox(s, x + 0.35, 13.55, 6.7, 0.55, sc["name"], 10.5, "navy", True)
        textbox(s, x + 0.35, 14.5, 6.7, 0.35, "Accesso SQL presso Guber", 8, "gray")
    star(s)


def scenario_slide(prs, page, sc):
    s = slide(prs); header(s, f"Scenario {sc['num']}", sc["name"], page)
    title(s, SCENARIO_MESSAGES[sc["num"]])
    textbox(s, 0.8, 2.42, 31.2, 0.5, SCENARIO_DESCRIPTIONS[sc["num"]], 9, "gray")
    # Visual left
    box(s, 0.8, 3.5, 13.8, 10.0, "white", "teal", True)
    textbox(s, 1.2, 3.9, 13.0, 0.35, "ARCHITETTURA DI RIFERIMENTO", 8, "dark_teal", True, PP_ALIGN.CENTER)
    for i, label in enumerate(sc["arch"]):
        y = 5.0 + i * 2.35
        architecture_node(s, 3.2, y, 9.0, label, "", "pale_blue" if i == 1 else "white")
        if i < 2:
            connector(s, 7.7, y + 1.45, 7.7, y + 2.25)
    # Content right
    textbox(s, 16.0, 3.75, 7.5, 0.4, "PERCHÉ CONSIDERARLO", 9, "dark_teal", True)
    bullets(s, 16.0, 4.35, 15.7, sc["pros"], 9.3)
    textbox(s, 16.0, 8.15, 7.5, 0.4, "PUNTI DI ATTENZIONE", 9, "blue", True)
    bullets(s, 16.0, 8.75, 15.7, sc["attention"], 9.3)
    textbox(s, 16.0, 12.15, 6.0, 0.35, "CONDIZIONI", 8, "gray", True)
    textbox(s, 21.0, 12.12, 10.5, 0.6, " · ".join(sc["gates"]), 8.2, "text", False, PP_ALIGN.RIGHT)
    box(s, 16.0, 13.25, 15.5, 1.55, "pale", "pale", True)
    textbox(s, 16.5, 13.68, 6.0, 0.35, "TCO PIATTAFORMA · 3 ANNI", 8, "gray", True)
    textbox(s, 23.0, 13.45, 7.7, 0.55, sc["tco"], 16, "dark_teal", True, PP_ALIGN.RIGHT)
    textbox(s, 16.0, 15.35, 15.5, 0.4,
            f"Una tantum {sc['oneoff']}  |  piattaforma mensile {sc['monthly']}", 8, "gray", False, PP_ALIGN.RIGHT)
    star(s)


def operating_model(prs):
    s = slide(prs); header(s, "NOVA", "Modello operativo", 11)
    title(s, "Account e chiavi possono restare Guber anche quando il run è affidato a Novigo")
    textbox(s, 0.9, 2.45, 31.0, 0.45,
            "Il modello operativo è indipendente dalla collocazione e deve coprire l'intero servizio NOVA.", 9, "gray")
    # table-like comparison
    headers = [("AM NOVIGO END-TO-END", "teal"), ("GUBER / TERZA PARTE", "blue")]
    for i, (h, col) in enumerate(headers):
        x = 8.6 + i * 11.7
        box(s, x, 3.8, 10.7, 0.85, col, col)
        textbox(s, x + 0.2, 4.05, 10.3, 0.3, h, 9, "white", True, PP_ALIGN.CENTER)
    rows = [
        ("Presidio", "Coordinato da un unico fornitore", "Distribuito tra owner e fornitori"),
        ("Perimetro", "Applicazione, DB, infrastruttura e orchestratore", "Da ricomporre nel contratto operativo"),
        ("Operations", "Monitor, backup, patching, incident e release", "Competenze e copertura da garantire"),
        ("Governance", "Account e approvazioni Guber", "Account e approvazioni Guber"),
        ("Handover", "Exit plan contrattuale", "Processo interno o di sourcing"),
    ]
    y = 4.65
    for idx, (label, left, right) in enumerate(rows):
        fill = "pale_blue" if idx % 2 else "white"
        box(s, 1.1, y, 7.2, 1.65, fill, "line")
        box(s, 8.6, y, 10.7, 1.65, fill, "line")
        box(s, 20.3, y, 10.7, 1.65, fill, "line")
        textbox(s, 1.45, y + 0.55, 6.5, 0.4, label, 9, "navy", True)
        textbox(s, 8.95, y + 0.45, 10.0, 0.75, left, 8.7, "text", False, PP_ALIGN.CENTER)
        textbox(s, 20.65, y + 0.45, 10.0, 0.75, right, 8.7, "text", False, PP_ALIGN.CENTER)
        y += 1.65
    star(s)


def economics_slide(prs):
    s = slide(prs); header(s, "NOVA", "Economics", 12)
    title(s, "I range TCO rendono confrontabili le alternative su un orizzonte di tre anni")
    textbox(s, 0.9, 2.45, 31.0, 0.45,
            "TCO piattaforma: costi una tantum più 36 mesi di costi ricorrenti.", 9, "gray")
    for i, sc in enumerate(SCENARIOS):
        col, row = i % 3, i // 3
        x, y = 0.9 + col * 10.7, 3.55 + row * 5.25
        box(s, x, y, 9.8, 4.55, "white", "line", True)
        box(s, x, y, 9.8, 0.65, "teal" if row == 0 else "blue", "teal" if row == 0 else "blue")
        textbox(s, x + 0.35, y + 0.9, 9.1, 0.65, sc["name"], 10, "navy", True)
        textbox(s, x + 0.35, y + 1.85, 4.4, 0.32, "UNA TANTUM", 7.5, "gray", True)
        textbox(s, x + 5.0, y + 1.85, 4.4, 0.32, "MENSILE", 7.5, "gray", True, PP_ALIGN.RIGHT)
        textbox(s, x + 0.35, y + 2.35, 4.4, 0.4, sc["oneoff"], 10.5, "text", True)
        textbox(s, x + 5.0, y + 2.35, 4.4, 0.4, sc["monthly"], 10.5, "text", True, PP_ALIGN.RIGHT)
        connector(s, x + 0.35, y + 3.0, x + 9.45, y + 3.0, "line", 0.8)
        textbox(s, x + 0.35, y + 3.35, 4.2, 0.35, "TCO 3 ANNI", 8, "dark_teal", True)
        textbox(s, x + 4.8, y + 3.2, 4.6, 0.5, sc["tco"], 14, "dark_teal", True, PP_ALIGN.RIGHT)
    textbox(s, 1.0, 14.75, 30.7, 0.8,
            "AM Novigo, team sostitutivo, licenze non note e disaster recovery geografico non inclusi.", 8.3, "gray", False, PP_ALIGN.CENTER)
    star(s)


def comparison_slide(prs):
    s = slide(prs); header(s, "NOVA", "Confronto delle alternative", 13)
    title(s, "Le alternative cambiano il punto di equilibrio tra velocità, controllo e complessità")
    criteria = ["Rapidità", "Rischio\nmigrazione", "Ownership\nGuber", "Accesso\nSQL", "Elasticità", "Semplicità\noperativa", "Reversibilità"]
    labels = {1: "Bassa", 2: "Medio-\nbassa", 3: "Media", 4: "Medio-\nalta", 5: "Alta"}
    x0, y0 = 0.8, 3.5
    widths = [9.4] + [3.25] * 7
    # Header
    x = x0
    for i, label in enumerate(["Scenario"] + criteria):
        box(s, x, y0, widths[i], 1.45, "navy", "white")
        textbox(s, x + 0.15, y0 + 0.4, widths[i] - 0.3, 0.7, label, 8.2, "white", True,
                PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        x += widths[i]
    palette = {1: "white", 2: "pale_blue", 3: "pale", 4: "pale", 5: "teal"}
    y = y0 + 1.45
    for r, sc in enumerate(SCENARIOS):
        x = x0
        box(s, x, y, widths[0], 1.55, "pale_blue", "white")
        textbox(s, x + 0.2, y + 0.5, widths[0] - 0.4, 0.4, sc["name"], 8.5, "navy", True)
        x += widths[0]
        for score in sc["scores"]:
            fill = palette[score]
            box(s, x, y, 3.25, 1.55, fill, "white")
            textbox(s, x + 0.1, y + 0.42, 3.05, 0.65, labels[score], 7.5,
                    "navy" if score != 5 else "white", score >= 4, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            x += 3.25
        y += 1.55
    textbox(s, 1.0, 14.2, 31.0, 0.5,
            "Valutazione qualitativa basata sulle informazioni disponibili.", 8, "gray", False, PP_ALIGN.CENTER)
    star(s)


def synthesis_slide(prs):
    s = slide(prs); header(s, "NOVA", "Sintesi", 14)
    title(s, "Le alternative rispondono a tre orientamenti infrastrutturali distinti")
    groups = [
        ("CONTINUITÀ", "Scenario 1", "Mantenere il primario nell'ambiente corrente", "teal"),
        ("TRASFERIMENTO DEL PRIMARIO", "Scenari 2–4", "Azure, AWS o infrastruttura on-premise Guber", "blue"),
        ("ACCESSO DATI CON REPLICA", "Scenari 5–6", "Replica read-only on-premise o Azure Guber", "navy"),
    ]
    for i, (label, ref, body, col) in enumerate(groups):
        x = 0.9 + i * 10.75
        box(s, x, 3.6, 9.8, 5.5, "white", col, True)
        box(s, x, 3.6, 9.8, 0.65, col, col)
        textbox(s, x + 0.35, 4.75, 9.1, 0.4, label, 8.5, "gray", True, PP_ALIGN.CENTER)
        textbox(s, x + 0.35, 5.65, 9.1, 0.5, ref, 14, "navy", True, PP_ALIGN.CENTER)
        textbox(s, x + 0.65, 6.75, 8.5, 1.3, body, 9.5, "text", False, PP_ALIGN.CENTER)
    textbox(s, 1.0, 10.7, 31.0, 0.45, "INFORMAZIONI CHE MODIFICANO IL CONFRONTO", 9, "dark_teal", True)
    items = [
        "Motore e versione DB", "Sizing e crescita", "SLA / RPO / RTO",
        "Standard cloud Guber", "Capacità on-premise", "Connettività", "Licensing",
    ]
    for i, item in enumerate(items):
        col, row = i % 4, i // 4
        x, y = 1.0 + col * 7.75, 11.55 + row * 1.65
        box(s, x, y, 7.0, 1.05, "pale_blue", "line", True)
        textbox(s, x + 0.2, y + 0.33, 6.6, 0.35, item, 8.5, "navy", True, PP_ALIGN.CENTER)
    star(s)


def appendix_cover(prs):
    s = slide(prs)
    box(s, 0, 0, 33.867, 19.05, "navy", "navy")
    s.shapes.add_picture(str(LOGO), Cm(0.65), Cm(0.55), width=Cm(3.2))
    s.shapes.add_picture(str(WAVE), Cm(22.7), Cm(0), height=Cm(19.05))
    textbox(s, 0.85, 10.8, 20.5, 0.4, "APPENDICE", 10, "teal", True)
    textbox(s, 0.85, 11.7, 22.0, 1.8, "Assunzioni, economics e\napprofondimenti tecnici", 24, "white", True)


def three_column_slide(prs, page, section, topic, slide_title, columns):
    s = slide(prs); header(s, section, topic, page); title(s, slide_title)
    for i, (label, values, col) in enumerate(columns):
        x = 0.9 + i * 10.7
        box(s, x, 3.3, 9.8, 10.5, "white", "line", True)
        box(s, x, 3.3, 9.8, 0.65, col, col)
        textbox(s, x + 0.35, 4.35, 9.1, 0.45, label, 9, "navy", True)
        bullets(s, x + 0.35, 5.25, 9.0, values, 9)
    star(s)


def assumptions_slide(prs):
    s = slide(prs); header(s, "Appendice", "Assunzioni economiche", 17)
    title(s, "Il confronto usa un carico di riferimento comune")
    rows = [
        ("Produzione", "4–8 vCPU · 16–32 GB RAM · 24x7 · alta disponibilità"),
        ("Database", "500 GB iniziali · crescita 20% annuo"),
        ("Non produzione", "1 ambiente al 50% della produzione"),
        ("Backup", "30 giorni · restore periodicamente testato"),
        ("Documentale", "1 TB iniziale"),
        ("Accesso SQL", "20 utenti autorizzati · massimo 5 concorrenti"),
        ("TCO", "Una tantum + 36 mesi di piattaforma"),
    ]
    y = 3.4
    for i, (label, value) in enumerate(rows):
        fill = "pale_blue" if i % 2 else "white"
        box(s, 1.0, y, 7.0, 1.45, fill, "white")
        box(s, 8.0, y, 15.0, 1.45, fill, "white")
        textbox(s, 1.35, y + 0.48, 6.3, 0.35, label, 9, "navy", True)
        textbox(s, 8.35, y + 0.42, 14.3, 0.5, value, 8.8, "text")
        y += 1.45
    box(s, 24.0, 3.4, 8.0, 10.15, "white", "line", True)
    box(s, 24.0, 3.4, 8.0, 0.65, "blue", "blue")
    textbox(s, 24.35, 4.45, 7.3, 0.4, "ESCLUSIONI", 9, "navy", True)
    bullets(s, 24.35, 5.25, 7.2, [
        "AM Novigo",
        "Team sostitutivo",
        "Licenze non note",
        "DR geografico",
        "Costi interni Guber",
        "IVA e sconti",
    ], 9)
    star(s)


def cost_detail_slide(prs, page, sc):
    s = slide(prs); header(s, "Appendice", f"Scenario {sc['num']} — economics", page)
    title(s, f"{sc['name']} — componenti del range")
    oneoff = {
        "01": ["Assessment e hardening", "Viste / accesso SQL", "Audit e test", "Documentazione"],
        "02": ["Assessment PaaS / IaaS", "Landing zone e rete", "Migrazione e test", "Cutover e rollback"],
        "03": ["Assessment RDS / EC2", "Landing zone e VPC", "Migrazione e test", "Cutover e rollback"],
        "04": ["Assessment e capacity", "Provisioning e licensing", "Migrazione e test", "HA, backup e cutover"],
        "05": ["Assessment replica", "Provisioning on-premise", "Setup e riconciliazione", "Accessi e monitor"],
        "06": ["Assessment replica", "Subscription e rete", "Setup e riconciliazione", "Accessi e monitor"],
    }[sc["num"]]
    recurring = {
        "01": ["Capacità incrementale", "Backup e logging", "Monitoring", "Networking"],
        "02": ["Compute / database", "Storage e backup", "Logging e security", "Networking"],
        "03": ["RDS / EC2", "Storage e backup", "CloudWatch / security", "Networking"],
        "04": ["Capacità e manutenzione", "Backup", "Monitoring", "Data center"],
        "05": ["Compute replica", "Storage e backup", "Monitor del lag", "Connettività"],
        "06": ["Database replica", "Storage e backup", "Monitor e audit", "Rete / egress"],
    }[sc["num"]]
    three = [
        ("UNA TANTUM", oneoff, "teal"),
        ("RICORRENTE", recurring, "blue"),
        ("VARIABILI PRINCIPALI", ["Motore e compatibilità", "HA / DR e SLA", "Volumi e IOPS", "Licenze e contratti"], "navy"),
    ]
    for i, (label, vals, col) in enumerate(three):
        x = 0.9 + i * 10.7
        box(s, x, 3.4, 9.8, 8.9, "white", "line", True)
        box(s, x, 3.4, 9.8, 0.65, col, col)
        textbox(s, x + 0.35, 4.45, 9.1, 0.4, label, 9, "navy", True)
        bullets(s, x + 0.35, 5.25, 9.0, vals, 9.2)
    box(s, 1.0, 13.25, 30.9, 1.6, "pale", "pale", True)
    values = [("Una tantum", sc["oneoff"]), ("Mensile", sc["monthly"]), ("TCO 3 anni", sc["tco"])]
    for i, (label, value) in enumerate(values):
        x = 1.5 + i * 10.0
        textbox(s, x, 13.62, 4.4, 0.35, label.upper(), 7.5, "gray", True)
        textbox(s, x + 4.3, 13.45, 5.0, 0.55, value, 12.5, "dark_teal", True, PP_ALIGN.RIGHT)
    star(s)


def security_slide(prs):
    three_column_slide(prs, 24, "Appendice", "Sicurezza e governance",
                       "Ownership e responsabilità devono restare esplicite",
                       [
                           ("GUBER", ["Account e infrastruttura", "Key ownership e policy", "Approvazioni e audit"], "teal"),
                           ("NOVIGO O ALTRO OPERATORE", ["Run operativo delegato", "Backup, patching e monitor", "Incident e release"], "blue"),
                           ("CONTROLLO COMUNE", ["Accesso privato e MFA", "Segregazione dei ruoli", "Logging, revoca ed exit plan"], "navy"),
                       ])


def storage_slide(prs):
    three_column_slide(prs, 25, "Appendice", "Storage documentale",
                       "La tecnologia dipende dal requisito di storage e dalla collocazione",
                       [
                           ("AWS S3", ["Lifecycle e versioning", "KMS e policy", "Egress da valutare"], "teal"),
                           ("AZURE STORAGE", ["Identity e policy Azure", "Tier e retention", "Connettività privata"], "blue"),
                           ("OBJECT STORAGE ON-PREMISE", ["Controllo locale", "Capacity e operations", "Durabilità da progettare"], "navy"),
                       ])


def orchestrator_slide(prs):
    three_column_slide(prs, 26, "Appendice", "Orchestratore e rete",
                       "La separazione dei componenti modifica dipendenze e responsabilità",
                       [
                           ("DB-ONLY", ["Database o replica spostati", "Applicazione e orchestratore invariati", "Dipendenza dalla WAN"], "teal"),
                           ("PERCORSO PER FASI", ["Replica o DB come primo passo", "Verifica prima di estendere lo scope", "Transitorio da governare"], "blue"),
                           ("FULL-STACK", ["Applicazione, DB e orchestratore", "Target più coerente", "Scope e cutover maggiori"], "navy"),
                       ])


def paas_slide(prs):
    s = slide(prs); header(s, "Appendice", "Deployment database", 27)
    title(s, "La compatibilità del motore determina il livello di servizio gestito")
    cols = [
        ("SERVIZIO GESTITO · PAAS / RDS", [
            "Patching e backup maggiormente gestiti",
            "Alta disponibilità standardizzata",
            "Minore carico operativo",
            "Compatibilità e feature da verificare",
        ], "teal"),
        ("VM / EC2 · SELF-MANAGED", [
            "Maggiore controllo e compatibilità",
            "Responsabilità DBA e sistema operativo",
            "Backup, HA e patching da costruire",
            "Licensing e TCO operativo da verificare",
        ], "blue"),
    ]
    for i, (label, vals, col) in enumerate(cols):
        x = 1.0 + i * 16.0
        box(s, x, 3.6, 14.9, 9.7, "white", col, True)
        box(s, x, 3.6, 14.9, 0.65, col, col)
        textbox(s, x + 0.5, 4.75, 13.9, 0.45, label, 10, "navy", True, PP_ALIGN.CENTER)
        bullets(s, x + 0.8, 6.0, 13.3, vals, 10)
    star(s)


def raci_slide(prs):
    s = slide(prs); header(s, "Appendice", "Modello operativo", 28)
    title(s, "Il run deve coprire applicazione, database e ambiente")
    rows = [
        ("Applicazione e release", "Novigo esegue", "Owner designato esegue"),
        ("Database e performance", "Novigo esegue", "DBA Guber / terzo"),
        ("Infrastruttura / cloud", "Novigo opera con delega", "Guber / terzo"),
        ("Backup e restore", "Novigo esegue e testa", "Guber / terzo"),
        ("Security operations", "Novigo opera; Guber governa", "Guber / terzo"),
        ("Account e chiavi", "Guber governa", "Guber governa"),
        ("Incident e problem", "Presidio coordinato", "Coordinamento multi-owner"),
    ]
    headers = ["Ambito", "AM Novigo end-to-end", "Guber / terza parte"]
    widths = [9.0, 11.5, 11.5]
    x0, y0 = 0.9, 3.4
    x = x0
    for i, h in enumerate(headers):
        box(s, x, y0, widths[i], 1.0, "navy" if i == 0 else ("teal" if i == 1 else "blue"), "white")
        textbox(s, x + 0.2, y0 + 0.3, widths[i] - 0.4, 0.4, h, 8.5, "white", True,
                PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        x += widths[i]
    y = y0 + 1.0
    for r, row in enumerate(rows):
        x = x0
        fill = "pale_blue" if r % 2 else "white"
        for i, value in enumerate(row):
            box(s, x, y, widths[i], 1.5, fill, "white")
            textbox(s, x + 0.25, y + 0.45, widths[i] - 0.5, 0.55, value, 8.6,
                    "navy" if i == 0 else "text", i == 0, PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
            x += widths[i]
        y += 1.5
    star(s)


def path_slide(prs):
    s = slide(prs); header(s, "Appendice", "Percorso di approfondimento", 29)
    title(s, "Quattro approfondimenti riducono l'incertezza tecnica ed economica")
    phases = [
        ("1", "Assessment", "Motore, sizing, licenze e compatibilità"),
        ("2", "Economics", "Contratti, TCO e modello operativo"),
        ("3", "Prova tecnica", "Replica o migrazione pilota"),
        ("4", "Piano di transizione", "Cutover, rollback, RACI e SLA"),
    ]
    connector(s, 3.2, 8.4, 29.8, 8.4, "teal", 2)
    for i, (n, name, body) in enumerate(phases):
        x = 1.1 + i * 8.0
        box(s, x, 4.4, 7.0, 7.7, "white", "line", True)
        box(s, x + 0.35, 4.85, 1.0, 0.72, "teal", "teal", True)
        textbox(s, x + 0.35, 5.05, 1.0, 0.3, n, 8.5, "white", True, PP_ALIGN.CENTER)
        textbox(s, x + 0.4, 6.3, 6.2, 0.5, name, 11, "navy", True)
        textbox(s, x + 0.4, 7.35, 6.2, 1.7, body, 8.7, "text")
    textbox(s, 1.0, 13.7, 31.0, 0.45,
            "Sequenza indicativa, attivabile sulle alternative che Guber sceglierà di approfondire.", 8.5, "gray", False, PP_ALIGN.CENTER)
    star(s)


def build():
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    clean_cover(prs)
    context_slide(prs)
    as_is_slide(prs)
    alternatives_slide(prs)
    page = 5
    for sc in SCENARIOS:
        scenario_slide(prs, page, sc)
        page += 1
    operating_model(prs)
    economics_slide(prs)
    comparison_slide(prs)
    synthesis_slide(prs)
    appendix_cover(prs)
    three_column_slide(prs, 16, "Appendice", "Evidenze e informazioni mancanti",
                       "Fatti, deduzioni e dati mancanti restano separati",
                       [
                           ("FATTI CONFERMATI", ["Accesso SQL richiesto", "Replica read-only accettata", "Azure e on-premise in agenda", "AWS confermato dall'utente", "AM Novigo end-to-end"], "teal"),
                           ("DEDUZIONI", ["Migrare il primario non è l'unico modo", "La replica riduce il rischio sul transazionale", "La rete diventa parte dell'architettura", "Key ownership è una scelta di governance"], "blue"),
                           ("DATI MANCANTI", ["Motore e licensing", "Sizing, crescita e IOPS", "SLA, RPO / RTO", "Landing zone e connettività", "Capacità on-premise e contratti"], "navy"),
                       ])
    assumptions_slide(prs)
    page = 18
    for sc in SCENARIOS:
        cost_detail_slide(prs, page, sc)
        page += 1
    security_slide(prs)
    storage_slide(prs)
    orchestrator_slide(prs)
    paas_slide(prs)
    raci_slide(prs)
    path_slide(prs)
    props = prs.core_properties
    props.title = "NOVA — alternative infrastrutturali e modelli operativi"
    props.subject = "Decision pack Guber — revisione visuale v2"
    props.author = "TXT e-solutions / Novigo"
    prs.save(OUT)
    print(OUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
