from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
REPO = ROOT.parent
TEMPLATE = REPO / "docs" / "template.pptx"
ASSETS = REPO / "2026-06-30-data-pipeline" / "generated-assets"
LOGO = ASSETS / "novigo-header-logo-template.png"
WAVE = ASSETS / "template-cover-wave.png"
STAR = ASSETS / "txt-star-mark-template.png"
OUT = ROOT / "attempts" / "NOVA Guber - visual prototype v2.pptx"

FONT = "Poppins"
C = {
    "navy": RGBColor(0x00, 0x44, 0x81),
    "blue": RGBColor(0x2A, 0x86, 0xCA),
    "sky": RGBColor(0x5B, 0xBE, 0xFE),
    "teal": RGBColor(0x2D, 0xCC, 0xCD),
    "dark_teal": RGBColor(0x00, 0x8E, 0x96),
    "pale": RGBColor(0xEC, 0xF8, 0xF8),
    "pale_blue": RGBColor(0xEF, 0xF6, 0xFB),
    "text": RGBColor(0x21, 0x25, 0x29),
    "gray": RGBColor(0x86, 0x8E, 0x95),
    "line": RGBColor(0xB9, 0xDA, 0xE5),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}


def clear(prs):
    ids = prs.slides._sldIdLst
    for slide_id in list(ids):
        prs.part.drop_rel(slide_id.rId)
        ids.remove(slide_id)


def slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(s, x, y, w, h, value="", size=10, color="text", bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.02)
    tf.margin_top = tf.margin_bottom = Cm(0.01)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = C[color]
    return box


def box(s, x, y, w, h, fill="white", line="line", rounded=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = s.shapes.add_shape(kind, Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C[fill]
    shp.line.color.rgb = C[line]
    shp.line.width = Pt(0.8)
    return shp


def connector(s, x1, y1, x2, y2, color="teal", width=1.2):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    c.line.color.rgb = C[color]
    c.line.width = Pt(width)


def header(s, section, topic, page):
    s.shapes.add_picture(str(LOGO), Cm(0.38), Cm(0.24), width=Cm(2.55))
    label = s.shapes.add_textbox(Cm(3.25), Cm(0.37), Cm(20), Cm(0.38))
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{section} | "; r.font.name = FONT; r.font.size = Pt(8); r.font.color.rgb = C["gray"]
    r = p.add_run(); r.text = topic; r.font.name = FONT; r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = C["text"]
    textbox(s, 31.9, 0.36, 0.7, 0.35, str(page), 8, "gray", align=PP_ALIGN.RIGHT)
    box(s, 0.38, 0.93, 1.35, 0.055, "teal", "teal")
    box(s, 32.5, 0.93, 0.75, 0.055, "blue", "blue")


def title(s, value):
    textbox(s, 0.75, 1.35, 31.5, 0.85, value, 18, "text", True)


def bullets(s, x, y, w, values, size=9):
    b = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(5.0))
    tf = b.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.02)
    for i, value in enumerate(values):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {value}"
        p.space_after = Pt(5)
        p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = C["text"]


def cover(prs):
    s = slide(prs)
    s.shapes.add_picture(str(LOGO), Cm(0.65), Cm(0.55), width=Cm(3.2))
    s.shapes.add_picture(str(WAVE), Cm(22.7), Cm(0), height=Cm(19.05))
    textbox(s, 0.75, 12.8, 21.0, 0.7, "NOVA", 19, "teal", True)
    textbox(s, 0.75, 13.55, 21.0, 1.25, "Alternative infrastrutturali\ne modelli operativi", 22, "text", True)
    textbox(s, 0.75, 15.2, 16.0, 0.4, "Guber", 10, "gray")
    textbox(s, 0.75, 17.5, 8.0, 0.35, "31 luglio 2026", 8, "gray")


def context(prs):
    s = slide(prs); header(s, "NOVA", "Contesto, esigenza e obiettivi", 2)
    xline = 5.2
    connector(s, xline, 3.0, xline, 15.6, "teal", 1.1)
    rows = [
        ("Contesto", "Guber richiede un accesso SQL governato ai dati NOVA",
         "L'esigenza riguarda verifiche ed estrazioni, mantenendo continuità operativa e tracciabilità degli accessi."),
        ("Esigenza", "Separare l'autonomia sui dati dal rischio sul transazionale",
         "La collocazione deve bilanciare controllo Guber, sicurezza, responsabilità operative e impatto della transizione."),
        ("Obiettivi", "Confrontare sei alternative su criteri omogenei",
         "Il confronto considera collocazione, accesso SQL, modello operativo, rischio, reversibilità ed economics."),
    ]
    for i, (label, lead, body) in enumerate(rows):
        y = 3.0 + i * 4.35
        textbox(s, 1.1, y + 0.3, 3.6, 0.4, label, 10, "dark_teal", True, PP_ALIGN.RIGHT)
        textbox(s, 6.2, y, 24.7, 0.55, lead, 12, "text", True)
        textbox(s, 6.2, y + 0.85, 24.7, 1.25, body, 9.5, "text")
    s.shapes.add_picture(str(STAR), Cm(31.9), Cm(17.3), width=Cm(0.65))


def scenario(prs):
    s = slide(prs); header(s, "Scenario 1", "Infrastruttura attuale", 3)
    title(s, "La continuità riduce il rischio di transizione, ma limita il controllo infrastrutturale di Guber")
    textbox(s, 1.0, 2.45, 30.5, 0.5,
            "Il database primario resta nell'ambiente corrente; l'accesso SQL viene segregato e tracciato.", 9, "gray")
    # Left visual
    box(s, 1.0, 3.6, 13.5, 9.4, "white", "teal", True)
    textbox(s, 1.5, 4.0, 12.5, 0.4, "ARCHITETTURA DI RIFERIMENTO", 8, "dark_teal", True, PP_ALIGN.CENTER)
    for i, (label, sub) in enumerate([
        ("Utenti Guber", "SQL read-only"),
        ("Accesso governato", "ruoli · viste · audit"),
        ("NOVA attuale", "database primario"),
    ]):
        y = 5.1 + i * 2.4
        box(s, 3.4, y, 8.7, 1.35, "pale_blue" if i == 1 else "white", "line", True)
        textbox(s, 3.7, y + 0.28, 8.1, 0.35, label, 10, "navy", True, PP_ALIGN.CENTER)
        textbox(s, 3.7, y + 0.78, 8.1, 0.28, sub, 7.5, "gray", False, PP_ALIGN.CENTER)
        if i < 2:
            connector(s, 7.75, y + 1.35, 7.75, y + 2.35)
    # Right content
    textbox(s, 16.0, 3.75, 7.0, 0.4, "PERCHÉ CONSIDERARLO", 9, "dark_teal", True)
    bullets(s, 16.0, 4.4, 15.7, [
        "Continuità del servizio e tempi di attivazione più brevi",
        "Nessuna migrazione immediata del database primario",
        "Alta reversibilità della scelta",
    ], 9.5)
    textbox(s, 16.0, 8.55, 7.0, 0.4, "PUNTI DI ATTENZIONE", 9, "blue", True)
    bullets(s, 16.0, 9.2, 15.7, [
        "Ownership infrastrutturale Guber limitata",
        "Accesso SQL da isolare dal workload applicativo",
        "SLA e capacità dell'ambiente corrente da confermare",
    ], 9.5)
    box(s, 16.0, 13.3, 15.5, 1.5, "pale", "pale", True)
    textbox(s, 16.5, 13.65, 5.5, 0.35, "TCO PIATTAFORMA · 3 ANNI", 8, "gray", True)
    textbox(s, 23.0, 13.45, 7.7, 0.55, "€ 40k–124k", 16, "dark_teal", True, PP_ALIGN.RIGHT)
    s.shapes.add_picture(str(STAR), Cm(31.9), Cm(17.3), width=Cm(0.65))


def main():
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    cover(prs); context(prs); scenario(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
