from pathlib import Path
from posixpath import normpath
import re
import sys
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PPTX = (
    Path(sys.argv[1]).resolve()
    if len(sys.argv) > 1
    else ROOT / "2026_GUBER_001 - Guber - NOVA Scenari infrastrutturali.pptx"
)


def fail(message):
    print(f"FAIL: {message}")
    sys.exit(1)


if not PPTX.exists() or PPTX.stat().st_size < 100_000:
    fail("PPTX missing or unexpectedly small")

prs = Presentation(PPTX)
if len(prs.slides) != 29:
    fail(f"expected 29 slides, found {len(prs.slides)}")

titles = []
all_text = []
for idx, slide in enumerate(prs.slides, start=1):
    slide_text = []
    for shp in slide.shapes:
        if hasattr(shp, "text") and shp.text.strip():
            slide_text.append(shp.text.strip())
    all_text.extend(slide_text)
    titles.append(slide_text[0] if slide_text else "")
    if idx not in (1, 15) and len(" ".join(slide_text)) < 40:
        fail(f"slide {idx} has too little text")

joined = "\n".join(all_text)
for required in [
    "Replica on-premise Guber",
    "Replica su Azure Guber",
    "AWS eu-central-1",
    "AM NOVIGO END-TO-END",
    "€ 108k–315k",
]:
    if required not in joined:
        fail(f"missing required content: {required}")

with zipfile.ZipFile(PPTX) as zf:
    bad = zf.testzip()
    if bad:
        fail(f"corrupt ZIP member: {bad}")
    names = set(zf.namelist())
    for name in names:
        if name.endswith((".xml", ".rels")):
            try:
                ET.fromstring(zf.read(name))
            except ET.ParseError as exc:
                fail(f"invalid XML in {name}: {exc}")
        if name.endswith(".rels"):
            root = ET.fromstring(zf.read(name))
            if name == "_rels/.rels":
                base = ""
            else:
                owner = name.replace("/_rels/", "/").removesuffix(".rels")
                base = str(Path(owner).parent).replace("\\", "/")
            for rel in root:
                if rel.attrib.get("TargetMode") == "External":
                    continue
                target = rel.attrib.get("Target", "")
                resolved = normpath(f"{base}/{target}").lstrip("/")
                if resolved not in names:
                    fail(f"broken relationship in {name}: {target} -> {resolved}")
    slide_xml = "\n".join(
        zf.read(name).decode("utf-8", errors="ignore")
        for name in names
        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
    )
    if re.search(r'<a:ext cx="-|<a:ext [^>]*cy="-', slide_xml):
        fail("negative shape extents detected")

print(f"OK: {PPTX.name}")
print(f"slides={len(prs.slides)} size={PPTX.stat().st_size}")
